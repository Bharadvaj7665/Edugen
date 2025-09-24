from fastapi import FastAPI, HTTPException, UploadFile, File, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse
from pydantic import BaseModel
from typing import List, Optional, Dict, Any
import os
import uuid
import asyncio
from datetime import datetime
import aiofiles
import json
import base64
from io import BytesIO

# AI Integration
from openai import AsyncOpenAI
from PIL import Image

# MongoDB
from motor.motor_asyncio import AsyncIOMotorClient
from pymongo import ASCENDING
import pymongo

# Document processing
import PyPDF2
from docx import Document as DocxDocument

# PPT Generation
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
import tempfile
import logging

# Load environment variables
from dotenv import load_dotenv
load_dotenv()

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = FastAPI(title="EduMind AI - Learning Assistant", version="1.0.0")

# CORS setup
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*","http://localhost:3000"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# MongoDB setup
MONGO_URL = os.environ.get('MONGO_URL', 'mongodb://localhost:27017/edumind_db')
client = AsyncIOMotorClient(MONGO_URL)
db = client.edumind_db

# Collections
chat_sessions_collection = db.chat_sessions
messages_collection = db.messages
documents_collection = db.documents
slides_collection = db.slides
slide_images_collection = db.slide_images
flashcards_collection = db.flashcards
flashcard_sets_collection = db.flashcard_sets
mcqs_collection = db.mcqs
mcq_sets_collection = db.mcq_sets
podcasts_collection = db.podcasts
podcast_sets_collection = db.podcast_sets
podcast_transcripts_collection = db.podcast_transcripts

# OpenAI API setup
OPENAI_API_KEY = os.environ.get('OPENAI_API_KEY')
if not OPENAI_API_KEY:
    raise ValueError("OPENAI_API_KEY environment variable is required")

# Configure OpenAI client
openai_client = AsyncOpenAI(api_key=OPENAI_API_KEY)

# File upload directory
UPLOAD_DIR = "/tmp/uploads"
EXPORT_DIR = "/tmp/exports"
IMAGES_DIR = "/tmp/slide_images"
os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(EXPORT_DIR, exist_ok=True)
os.makedirs(IMAGES_DIR, exist_ok=True)

# Pydantic models
class ChatMessage(BaseModel):
    session_id: str
    message: str
    file_id: Optional[str] = None

class ChatResponse(BaseModel):
    response: str
    session_id: str
    timestamp: datetime

class ChatSession(BaseModel):
    session_id: str
    title: str
    created_at: datetime
    updated_at: datetime

class DocumentInfo(BaseModel):
    file_id: str
    filename: str
    file_type: str
    uploaded_at: datetime

class SlideGenerationRequest(BaseModel):
    file_id: str
    session_id: str
    slide_count: int = 8
    include_images: bool = False  # New option for visual enhancement

class ImageGenerationRequest(BaseModel):
    slide_id: str
    prompt: Optional[str] = None  # Custom prompt or auto-generated from slide content

class TitleFormatting(BaseModel):
    font_size: Optional[int] = 24
    font_family: Optional[str] = "Calibri"
    text_color: Optional[str] = "#1f2937"
    is_bold: Optional[bool] = True
    is_italic: Optional[bool] = False

class BodyFormatting(BaseModel):
    font_size: Optional[int] = 18
    font_family: Optional[str] = "Calibri"
    text_color: Optional[str] = "#333333"
    bullet_style: Optional[str] = "bullet"
    is_bold: Optional[bool] = False
    is_italic: Optional[bool] = False

class SlideImage(BaseModel):
    image_id: str
    slide_id: str
    image_base64: str
    prompt_used: str
    generated_at: datetime

class Slide(BaseModel):
    slide_id: str
    slide_number: int
    title: str
    content: List[str]
    speaker_notes: str
    title_formatting: Optional[TitleFormatting] = None
    body_formatting: Optional[BodyFormatting] = None
    image_id: Optional[str] = None  # Reference to generated image

class SlideEditRequest(BaseModel):
    slide_id: str
    title: str
    content: List[str]
    speaker_notes: str
    title_formatting: Optional[TitleFormatting] = None
    body_formatting: Optional[BodyFormatting] = None

class PresentationTheme(BaseModel):
    theme_name: str
    background_color: str
    title_color: str
    text_color: str
    font_family: str
    font_size: int

class ExportRequest(BaseModel):
    file_id: str
    theme: PresentationTheme

class SlidePresentation(BaseModel):
    slides: List[Slide]
    total_slides: int
    document_title: str
    generated_at: datetime
    is_editable: bool = True

# Flashcard Models
class FlashcardGenerationRequest(BaseModel):
    file_id: str
    session_id: str
    card_count: int = 20
    card_type: str = "qa"  # qa, true_false, fill_blank
    difficulty: str = "mixed"  # easy, medium, hard, mixed
    topics: Optional[List[str]] = None  # specific topics to focus on

class Flashcard(BaseModel):
    card_id: str
    question: str
    answer: str
    card_type: str = "qa"  # qa, true_false, fill_blank
    topic: str
    tags: List[str] = []
    difficulty: str = "medium"  # easy, medium, hard
    source_doc: str
    file_id: str
    created_at: datetime
    last_reviewed: Optional[datetime] = None
    review_count: int = 0
    correct_count: int = 0

class FlashcardSet(BaseModel):
    set_id: str
    set_name: str
    file_id: str
    session_id: str
    document_title: str
    total_cards: int
    card_type: str = "qa"
    generated_at: datetime
    last_modified: datetime
    is_editable: bool = True

class FlashcardEditRequest(BaseModel):
    card_id: str
    question: str
    answer: str
    topic: str
    tags: List[str] = []
    difficulty: str = "medium"

class FlashcardExportRequest(BaseModel):
    set_id: str
    export_format: str = "csv"  # csv, anki, json

# MCQ Models
class MCQGenerationRequest(BaseModel):
    file_id: str
    session_id: str
    question_count: int = 15
    question_type: str = "single_correct"  # single_correct, multiple_correct, true_false
    difficulty: str = "mixed"  # easy, medium, hard, mixed
    topics: Optional[List[str]] = None  # specific topics to focus on

class MCQOption(BaseModel):
    option_text: str
    is_correct: bool

class MCQ(BaseModel):
    question_id: str
    question_text: str
    options: List[MCQOption]
    explanation: str
    question_type: str = "single_correct"  # single_correct, multiple_correct, true_false
    topic: str
    difficulty: str = "medium"  # easy, medium, hard
    bloom_level: str = "remember"  # remember, understand, apply, analyze
    source_doc: str
    file_id: str
    created_at: datetime
    last_reviewed: Optional[datetime] = None
    times_attempted: int = 0
    times_correct: int = 0

class MCQSet(BaseModel):
    set_id: str
    set_name: str
    file_id: str
    session_id: str
    document_title: str
    total_questions: int
    question_type: str = "single_correct"
    difficulty_distribution: Dict[str, int] = {}  # {"easy": 5, "medium": 8, "hard": 2}
    generated_at: datetime
    last_modified: datetime
    is_editable: bool = True

class MCQEditRequest(BaseModel):
    question_id: str
    question_text: str
    options: List[MCQOption]
    explanation: str
    topic: str
    difficulty: str = "medium"
    bloom_level: str = "remember"

class MCQExportRequest(BaseModel):
    set_id: str
    export_format: str = "csv"  # csv, pdf, json

# Podcast Models
class PodcastTranscriptRequest(BaseModel):
    file_id: str
    session_id: str
    podcast_length: str = "medium"  # quick, medium, comprehensive
    content_focus: str = "full_document"  # full_document, key_concepts, summary

class PodcastGenerationRequest(BaseModel):
    transcript_id: str
    voice_style: str = "neutral"  # neutral, enthusiastic, formal, conversational
    voice_gender: str = "female"  # male, female
    voice_accent: str = "american"  # american, british, indian, australian, canadian
    include_intro_outro: bool = True

class PodcastTranscript(BaseModel):
    transcript_id: str
    title: str
    description: str
    script_text: str
    podcast_length: str
    estimated_duration: str
    source_doc: str
    file_id: str
    created_at: datetime
    is_editable: bool = True

class PodcastChapter(BaseModel):
    title: str
    timestamp: str  # MM:SS format
    duration: int  # seconds

class Podcast(BaseModel):
    podcast_id: str
    transcript_id: str
    title: str
    description: str
    script_text: str
    audio_file_path: str
    duration_seconds: int
    duration_formatted: str  # MM:SS format
    voice_style: str
    voice_gender: str
    voice_accent: str
    language: str = "en"
    chapters: List[PodcastChapter] = []
    file_format: str = "mp3"  # mp3, wav
    file_size_bytes: int
    source_doc: str
    file_id: str
    created_at: datetime
    is_processing: bool = False
    processing_status: str = "completed"  # generating, processing, completed, failed

class PodcastSet(BaseModel):
    set_id: str
    set_name: str
    file_id: str
    session_id: str
    document_title: str
    total_podcasts: int = 1  # For now, one podcast per document
    generated_at: datetime
    last_modified: datetime
    is_editable: bool = True

class PodcastEditRequest(BaseModel):
    transcript_id: str
    title: str
    script_text: str

class PodcastRegenerateRequest(BaseModel):
    transcript_id: str
    voice_style: str = "neutral"
    voice_gender: str = "female"
    voice_accent: str = "american"
    include_intro_outro: bool = True

class PodcastExportRequest(BaseModel):
    podcast_id: str
    export_format: str = "mp3"  # mp3, wav

# Enhanced presentation themes with visual design elements
PRESENTATION_THEMES = {
    "professional": {
        "theme_name": "Professional",
        "background_color": "#FFFFFF",
        "title_color": "#2C3E50",
        "text_color": "#34495E",
        "font_family": "Calibri",
        "font_size": 18,
        "design_style": "clean_minimal"
    },
    "modern": {
        "theme_name": "Modern",
        "background_color": "#F8F9FA",
        "title_color": "#6366F1",
        "text_color": "#374151",
        "font_family": "Arial",
        "font_size": 18,
        "design_style": "gradient_accent"
    },
    "creative": {
        "theme_name": "Creative",
        "background_color": "#FEF7FF",
        "title_color": "#7C3AED",
        "text_color": "#4B5563",
        "font_family": "Segoe UI",
        "font_size": 18,
        "design_style": "artistic_vibrant"
    },
    "minimal": {
        "theme_name": "Minimal",
        "background_color": "#FAFAFA",
        "title_color": "#1F2937",
        "text_color": "#6B7280",
        "font_family": "Helvetica",
        "font_size": 18,
        "design_style": "ultra_clean"
    },
    "corporate": {
        "theme_name": "Corporate",
        "background_color": "#F1F5F9",
        "title_color": "#0F172A",
        "text_color": "#334155",
        "font_family": "Times New Roman",
        "font_size": 18,
        "design_style": "business_formal"
    }
}

# Helper functions
async def extract_text_from_file(file_path: str, file_type: str) -> str:
    """Extract text from uploaded files"""
    try:
        if file_type == "application/pdf":
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                text = ""
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
                return text
        elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            doc = DocxDocument(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        elif file_type == "text/plain":
            async with aiofiles.open(file_path, 'r', encoding='utf-8') as file:
                return await file.read()
        else:
            return ""
    except Exception as e:
        logger.error(f"Error extracting text: {e}")
        return ""

async def get_or_create_chat_session(session_id: str) -> dict:
    """Get existing chat session or create new one"""
    session = await chat_sessions_collection.find_one({"session_id": session_id})
    if not session:
        session_data = {
            "session_id": session_id,
            "title": "New Chat Session",
            "created_at": datetime.utcnow(),
            "updated_at": datetime.utcnow()
        }
        await chat_sessions_collection.insert_one(session_data)
        return session_data
    return session

async def generate_slide_image(slide_title: str, slide_content: List[str], custom_prompt: Optional[str] = None) -> Optional[str]:
    """Generate an image for a slide using OpenAI DALL-E"""
    try:
        # Generate smart prompt based on slide content
        if custom_prompt:
            prompt = custom_prompt
        else:
            content_text = " ".join(slide_content[:3])  # Use first 3 bullet points
            prompt = f"Create a professional, educational illustration for a presentation slide titled '{slide_title}'. The slide covers: {content_text}. Style: clean, modern, suitable for business presentations, high quality, photorealistic or professional illustration style. No text overlay needed."
        
        logger.info(f"Generating image with prompt: {prompt[:100]}...")
        
        try:
            # Generate image using OpenAI DALL-E
            response = await openai_client.images.generate(
                model="dall-e-3",
                prompt=prompt,
                size="1024x1024",
                quality="standard",
                response_format="b64_json",
                n=1
            )
            
            # Extract base64 image data
            if response.data and len(response.data) > 0:
                image_base64 = response.data[0].b64_json
                logger.info("Image generated successfully with OpenAI DALL-E")
                return image_base64
            else:
                logger.warning("No image data returned from OpenAI")
                return None
                
        except Exception as e:
            logger.warning(f"OpenAI DALL-E failed: {e}")
            # For now, return None to allow the system to continue without images
            logger.warning("Image generation not available - continuing without images")
            return None
        
    except Exception as e:
        logger.error(f"Error generating image: {e}")
        return None

async def generate_slide_content(file_path: str, file_type: str, document_title: str, slide_count: int, include_images: bool = False) -> List[Dict[str, Any]]:
    """Generate content-specific slide content using OpenAI with optional images"""
    try:
        # Read document content
        document_text = await extract_text_from_file(file_path, file_type)
        
        slide_prompt = f"""
        Please analyze this specific document and create a presentation with EXACTLY {slide_count} slides that are directly based on the content provided. 
        Each slide must contain information that is SPECIFICALLY from this document - do not add generic information.
        
        Return the response as a JSON object with the following exact structure:

        {{
            "slides": [
                {{
                    "slide_id": "slide_1",
                    "slide_number": 1,
                    "title": "Specific title based on document content",
                    "content": ["Specific point 1 from document", "Specific point 2 from document", "Specific point 3 from document", "Specific point 4 from document", "Specific point 5 from document"],
                    "speaker_notes": "Comprehensive and detailed speaker notes (minimum 150-200 words) that provide thorough context, background information, key explanations, examples from the document, connections to broader concepts, practical applications, and specific guidance for the presenter on how to explain each bullet point effectively. Include specific quotes, data, or examples from the document where relevant. Explain why these points are important and how they relate to the overall topic. Provide additional context that would help an audience understand the significance of each point.",
                    "title_formatting": {{
                        "font_size": 24,
                        "font_family": "Calibri",
                        "text_color": "#1f2937",
                        "is_bold": true,
                        "is_italic": false
                    }},
                    "body_formatting": {{
                        "font_size": 18,
                        "font_family": "Calibri",
                        "text_color": "#333333",
                        "bullet_style": "bullet",
                        "is_bold": false,
                        "is_italic": false
                    }}
                }}
            ]
        }}

        CRITICAL REQUIREMENTS:
        1. Create EXACTLY {slide_count} slides (no more, no less)
        2. Each slide must have EXACTLY 4-5 bullet points (no more, no less)
        3. All content must be SPECIFIC to the uploaded document - extract actual facts, data, concepts, or ideas
        4. Slide titles should be descriptive and specific to the content
        5. Speaker notes must be COMPREHENSIVE and DETAILED (minimum 150-200 words each)
        6. Speaker notes should include:
           - Detailed explanation of each bullet point
           - Context and background information
           - Specific examples or data from the document
           - Connections to broader concepts
           - Practical applications or implications
           - Guidance for effective presentation delivery
        7. First slide should introduce the document's main topic
        8. Last slide should summarize key takeaways from the document
        9. Middle slides should cover main themes, concepts, or sections from the document
        10. Include specific examples, data, or quotes from the document where relevant
        11. Ensure proper JSON formatting
        12. Distribute content evenly across all {slide_count} slides

        Document title: {document_title}
        Document content: {document_text[:4000]}...
        
        Analyze the document content carefully and create slides that provide comprehensive coverage of the material with detailed speaker guidance.
        """
        
        # Get AI response from OpenAI
        response = await openai_client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {
                    "role": "system",
                    "content": "You are an expert presentation designer and educational content creator. Your task is to analyze documents and create highly specific, content-focused presentation slides that extract and organize the exact information from the uploaded document. Focus on clarity, educational value, and content accuracy. Create comprehensive, detailed speaker notes that provide deep context and explanation for each slide."
                },
                {
                    "role": "user",
                    "content": slide_prompt
                }
            ],
            max_tokens=8000,
            temperature=0.7
        )
        
        ai_response = response.choices[0].message.content
        
        # Clean and parse JSON response
        json_start = ai_response.find('{')
        json_end = ai_response.rfind('}') + 1
        
        if json_start != -1 and json_end != 0:
            json_content = ai_response[json_start:json_end]
            slide_data = json.loads(json_content)
            
            if "slides" in slide_data and isinstance(slide_data["slides"], list):
                # Add unique slide IDs and ensure proper structure
                processed_slides = []
                for i, slide in enumerate(slide_data["slides"]):
                    default_title_formatting = {
                        "font_size": 24,
                        "font_family": "Calibri",
                        "text_color": "#1f2937",
                        "is_bold": True,
                        "is_italic": False
                    }
                    
                    default_body_formatting = {
                        "font_size": 18,
                        "font_family": "Calibri",
                        "text_color": "#333333",
                        "bullet_style": "bullet",
                        "is_bold": False,
                        "is_italic": False
                    }
                    
                    slide_id = f"slide_{i+1}_{uuid.uuid4().hex[:8]}"
                    processed_slide = {
                        "slide_id": slide_id,
                        "slide_number": i + 1,
                        "title": slide.get("title", f"Slide {i+1}"),
                        "content": slide.get("content", [])[:5],  # Ensure max 5 bullet points
                        "speaker_notes": slide.get("speaker_notes", ""),
                        "title_formatting": slide.get("title_formatting", default_title_formatting),
                        "body_formatting": slide.get("body_formatting", default_body_formatting),
                        "image_id": None  # Will be populated if images are generated
                    }
                    # Ensure we have exactly 4-5 bullet points
                    while len(processed_slide["content"]) < 4:
                        processed_slide["content"].append("Additional point to be added")
                    
                    # Generate image if requested
                    if include_images:
                        try:
                            image_base64 = await generate_slide_image(
                                processed_slide["title"], 
                                processed_slide["content"]
                            )
                            if image_base64:
                                # Store image in database
                                image_id = f"img_{slide_id}_{uuid.uuid4().hex[:8]}"
                                image_data = {
                                    "image_id": image_id,
                                    "slide_id": slide_id,
                                    "image_base64": image_base64,
                                    "prompt_used": f"Educational illustration for '{processed_slide['title']}'",
                                    "generated_at": datetime.utcnow()
                                }
                                await slide_images_collection.insert_one(image_data)
                                processed_slide["image_id"] = image_id
                                logger.info(f"Generated image for slide {i+1}")
                        except Exception as e:
                            logger.warning(f"Failed to generate image for slide {i+1}: {e}")
                    
                    processed_slides.append(processed_slide)
                
                return processed_slides[:slide_count]  # Ensure exact count
        
        # Fallback if JSON parsing fails
        return create_fallback_slides(document_title, slide_count, include_images)
        
    except Exception as e:
        logger.error(f"Error generating slides: {e}")
        return create_fallback_slides(document_title, slide_count, include_images)

def create_fallback_slides(document_title: str, slide_count: int, include_images: bool = False) -> List[Dict[str, Any]]:
    """Create fallback slides if AI generation fails"""
    default_title_formatting = {
        "font_size": 24,
        "font_family": "Calibri",
        "text_color": "#1f2937",
        "is_bold": True,
        "is_italic": False
    }
    
    default_body_formatting = {
        "font_size": 18,
        "font_family": "Calibri",
        "text_color": "#333333",
        "bullet_style": "bullet",
        "is_bold": False,
        "is_italic": False
    }
    
    slides = []
    
    # Create the requested number of slides
    for i in range(slide_count):
        slide_id = f"fallback_{i+1}_{uuid.uuid4().hex[:8]}"
        
        if i == 0:  # First slide
            slide = {
                "slide_id": slide_id,
                "slide_number": i + 1,
                "title": f"Introduction to {document_title}",
                "content": [
                    "Document overview and main topic",
                    "Key themes to be explored",
                    "Learning objectives for this presentation",
                    "Structure and organization of content",
                    "Expected outcomes from this analysis"
                ],
                "speaker_notes": "Welcome to this presentation about the uploaded document. This introduction slide sets the stage for understanding the main themes and concepts that will be covered. Begin by providing context about the document's origin and purpose. Explain the learning objectives and what the audience can expect to gain from this presentation. Highlight the key themes that will be explored in subsequent slides. Discuss the structure of the presentation and how the content is organized to facilitate understanding. Conclude by setting expectations for the outcomes and takeaways that will be provided throughout the session. This foundational slide establishes credibility and prepares the audience for deeper exploration of the material.",
                "title_formatting": default_title_formatting,
                "body_formatting": default_body_formatting,
                "image_id": None
            }
        elif i == slide_count - 1:  # Last slide
            slide = {
                "slide_id": slide_id,
                "slide_number": i + 1,
                "title": "Summary and Takeaways",
                "content": [
                    "Main findings and conclusions from the document",
                    "Key insights and important points",
                    "Practical applications and implications",
                    "Areas for further exploration",
                    "Final thoughts and recommendations"
                ],
                "speaker_notes": "This concluding slide summarizes the most important points from the document and presentation. Review the main findings and conclusions, emphasizing their significance and relevance. Discuss the key insights that emerged from the analysis and why they matter. Explain the practical applications and real-world implications of the content covered. Suggest areas for further exploration or study that would build upon this foundation. Provide final thoughts and recommendations for action or continued learning. End with a strong conclusion that ties together all the major themes and leaves the audience with clear takeaways they can apply. This slide should reinforce the value of the presentation and inspire further engagement with the topic.",
                "title_formatting": default_title_formatting,
                "body_formatting": default_body_formatting,
                "image_id": None
            }
        else:  # Middle slides
            slide = {
                "slide_id": slide_id,
                "slide_number": i + 1,
                "title": f"Key Topic {i}",
                "content": [
                    f"Primary concept {i} identified in the document",
                    f"Important definition and terminology for topic {i}",
                    f"Core principle related to concept {i}",
                    f"Relationship between topic {i} and other ideas",
                    f"Context and background for topic {i}"
                ],
                "speaker_notes": f"This slide covers key topic {i} found in the document. Begin by explaining the primary concept and its importance within the broader context of the document. Provide clear definitions and explain any technical terminology that the audience needs to understand. Discuss the core principles associated with this topic and why they matter. Explain how this concept relates to other ideas presented in the document, showing connections and dependencies. Provide background context that helps the audience understand the significance of this topic. Use specific examples from the document to illustrate your points and make the content more relatable and memorable. This detailed exploration ensures comprehensive understanding of the material.",
                "title_formatting": default_title_formatting,
                "body_formatting": default_body_formatting,
                "image_id": None
            }
        
        slides.append(slide)
    
    return slides

async def generate_flashcard_content(file_path: str, file_type: str, document_title: str, card_count: int, card_type: str = "qa") -> List[Dict[str, Any]]:
    """Generate flashcard content using OpenAI"""
    try:
        # Read document content
        document_text = await extract_text_from_file(file_path, file_type)
        
        flashcard_prompt = f"""
        Please analyze this document and create EXACTLY {card_count} flashcards in Q&A format based on the content.
        
        Return the response as a JSON object with the following exact structure:
        
        {{
            "flashcards": [
                {{
                    "card_id": "fc_1",
                    "question": "Clear, concise question testing a specific concept",
                    "answer": "Complete but concise answer with key information",
                    "topic": "Main topic/subject area this card covers",
                    "tags": ["tag1", "tag2", "tag3"],
                    "difficulty": "easy/medium/hard"
                }}
            ]
        }}
        
        CRITICAL REQUIREMENTS:
        1. Create EXACTLY {card_count} flashcards (no more, no less)
        2. Each question should test ONE specific concept, definition, or fact
        3. Questions should be clear and unambiguous
        4. Answers should be complete but concise (1-3 sentences)
        5. Cover different topics from the document evenly
        6. Include a variety of question types:
           - Definitions ("What is...?")
           - Explanations ("How does...?", "Why does...?")
           - Facts ("When did...?", "Where is...?")
           - Processes ("What happens when...?")
           - Comparisons ("What's the difference between...?")
        7. Assign appropriate difficulty levels:
           - Easy: Basic definitions and simple facts
           - Medium: Explanations and processes
           - Hard: Complex relationships and analysis
        8. Generate relevant tags (2-4 per card) from document content
        9. Ensure proper JSON formatting
        10. Extract content ONLY from the provided document
        
        Document title: {document_title}
        Document content: {document_text[:4000]}...
        
        Focus on creating flashcards that will help students learn and retain the key information from this document.
        """
        
        # Get AI response from OpenAI
        response = await openai_client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {
                    "role": "system",
                    "content": "You are an expert educational content creator specializing in creating effective flashcards for learning. Your task is to analyze documents and create high-quality Q&A flashcards that test key concepts, definitions, facts, and understanding. Focus on creating cards that promote active recall and spaced repetition learning."
                },
                {
                    "role": "user",
                    "content": flashcard_prompt
                }
            ],
            max_tokens=8000,
            temperature=0.7
        )
        
        ai_response = response.choices[0].message.content
        
        # Clean and parse JSON response
        json_start = ai_response.find('{')
        json_end = ai_response.rfind('}') + 1
        
        if json_start != -1 and json_end != 0:
            json_content = ai_response[json_start:json_end]
            flashcard_data = json.loads(json_content)
            
            if "flashcards" in flashcard_data and isinstance(flashcard_data["flashcards"], list):
                # Add unique card IDs and ensure proper structure
                processed_cards = []
                for i, card in enumerate(flashcard_data["flashcards"]):
                    card_id = f"fc_{i+1}_{uuid.uuid4().hex[:8]}"
                    processed_card = {
                        "card_id": card_id,
                        "question": card.get("question", f"Question {i+1}"),
                        "answer": card.get("answer", f"Answer {i+1}"),
                        "card_type": card_type,
                        "topic": card.get("topic", "General"),
                        "tags": card.get("tags", []),
                        "difficulty": card.get("difficulty", "medium").lower(),
                        "source_doc": document_title,
                        "created_at": datetime.utcnow(),
                        "last_reviewed": None,
                        "review_count": 0,
                        "correct_count": 0
                    }
                    processed_cards.append(processed_card)
                
                return processed_cards[:card_count]  # Ensure exact count
        
        # Fallback if JSON parsing fails
        return create_fallback_flashcards(document_title, card_count, card_type)
        
    except Exception as e:
        logger.error(f"Error generating flashcards: {e}")
        return create_fallback_flashcards(document_title, card_count, card_type)

def create_fallback_flashcards(document_title: str, card_count: int, card_type: str = "qa") -> List[Dict[str, Any]]:
    """Create fallback flashcards if AI generation fails"""
    fallback_cards = []
    
    for i in range(min(card_count, 10)):  # Limit fallback to 10 cards max
        card_id = f"fallback_{i+1}_{uuid.uuid4().hex[:8]}"
        
        fallback_card = {
            "card_id": card_id,
            "question": f"What is the main topic of {document_title}?",
            "answer": f"This is a fallback flashcard for {document_title}. Please review the document and create custom flashcards.",
            "card_type": card_type,
            "topic": "General",
            "tags": ["fallback", "review"],
            "difficulty": "easy",
            "source_doc": document_title,
            "created_at": datetime.utcnow(),
            "last_reviewed": None,
            "review_count": 0,
            "correct_count": 0
        }
        fallback_cards.append(fallback_card)
    
    return fallback_cards

async def generate_mcq_content(file_path: str, file_type: str, document_title: str, question_count: int, question_type: str = "single_correct", difficulty: str = "mixed") -> List[Dict[str, Any]]:
    """Generate MCQ content using OpenAI"""
    try:
        # Read document content
        document_text = await extract_text_from_file(file_path, file_type)
        
        # Determine difficulty instruction
        difficulty_instruction = ""
        if difficulty == "easy":
            difficulty_instruction = "Focus on basic recall and simple understanding questions."
        elif difficulty == "medium":
            difficulty_instruction = "Create questions that test understanding and application of concepts."
        elif difficulty == "hard":
            difficulty_instruction = "Generate challenging questions that require analysis and critical thinking."
        else:  # mixed
            difficulty_instruction = "Create a mix of difficulty levels: 40% easy (recall), 40% medium (understanding), 20% hard (analysis)."
        
        mcq_prompt = f"""
        Please analyze this document and create EXACTLY {question_count} high-quality single-answer multiple-choice questions based on the content.
        
        {difficulty_instruction}
        
        Return the response as a JSON object with the following exact structure:
        
        {{
            "mcqs": [
                {{
                    "question_id": "mcq_1",
                    "question_text": "Clear, specific question testing a key concept",
                    "options": [
                        {{
                            "option_text": "First option (could be correct)",
                            "is_correct": true
                        }},
                        {{
                            "option_text": "Second option (distractor)",
                            "is_correct": false
                        }},
                        {{
                            "option_text": "Third option (distractor)",
                            "is_correct": false
                        }},
                        {{
                            "option_text": "Fourth option (distractor)",
                            "is_correct": false
                        }}
                    ],
                    "explanation": "Clear explanation of why the correct answer is right and why others are wrong",
                    "topic": "Main topic/subject area this question covers",
                    "difficulty": "easy/medium/hard",
                    "bloom_level": "remember/understand/apply/analyze"
                }}
            ]
        }}
        
        CRITICAL REQUIREMENTS:
        1. Create EXACTLY {question_count} multiple choice questions (no more, no less)
        2. Each question must have EXACTLY 4 options with EXACTLY 1 correct answer
        3. Questions should test key concepts, definitions, processes, and relationships from the document
        4. Create plausible distractors that could seem correct to someone who hasn't studied the material carefully
        5. Ensure questions are clear, unambiguous, and grammatically correct
        6. Provide comprehensive explanations that teach the concept
        7. Cover different topics from the document evenly
        8. Include variety in question stems:
           - "Which of the following..." (most common)
           - "What is the primary..." 
           - "According to the document..."
           - "The main difference between X and Y is..."
           - "Which statement is correct about..."
        9. Assign appropriate Bloom's taxonomy levels:
           - Remember: Basic recall of facts and definitions
           - Understand: Explanation and interpretation
           - Apply: Using knowledge in new situations  
           - Analyze: Breaking down complex information
        10. Assign difficulty levels appropriately
        11. Ensure proper JSON formatting
        12. Extract content ONLY from the provided document
        
        Document title: {document_title}
        Document content: {document_text[:4000]}...
        
        Focus on creating assessment-quality questions that would be suitable for testing student knowledge and understanding of this content.
        """
        
        # Get AI response from OpenAI
        response = await openai_client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {
                    "role": "system",
                    "content": "You are an expert educational assessment creator specializing in generating high-quality multiple choice questions. Your task is to analyze documents and create assessment-quality MCQs that test conceptual understanding, application, and recall. Focus on creating questions that are clear, have plausible distractors, and promote learning."
                },
                {
                    "role": "user",
                    "content": mcq_prompt
                }
            ],
            max_tokens=8000,
            temperature=0.7
        )
        
        ai_response = response.choices[0].message.content
        
        # Clean and parse JSON response
        json_start = ai_response.find('{')
        json_end = ai_response.rfind('}') + 1
        
        if json_start != -1 and json_end != 0:
            json_content = ai_response[json_start:json_end]
            mcq_data = json.loads(json_content)
            
            if "mcqs" in mcq_data and isinstance(mcq_data["mcqs"], list):
                # Process and validate MCQs
                processed_mcqs = []
                for i, mcq in enumerate(mcq_data["mcqs"]):
                    question_id = f"mcq_{i+1}_{uuid.uuid4().hex[:8]}"
                    
                    # Ensure proper options structure
                    options = mcq.get("options", [])
                    if len(options) != 4:
                        # Create fallback options if not properly generated
                        options = [
                            {"option_text": f"Option A", "is_correct": True},
                            {"option_text": f"Option B", "is_correct": False},
                            {"option_text": f"Option C", "is_correct": False},
                            {"option_text": f"Option D", "is_correct": False}
                        ]
                    
                    # Validate exactly one correct answer
                    correct_count = sum(1 for opt in options if opt.get("is_correct", False))
                    if correct_count != 1:
                        # Fix: ensure exactly one correct answer
                        for opt in options:
                            opt["is_correct"] = False
                        options[0]["is_correct"] = True
                    
                    processed_mcq = {
                        "question_id": question_id,
                        "question_text": mcq.get("question_text", f"Question {i+1}"),
                        "options": options,
                        "explanation": mcq.get("explanation", "Explanation not provided."),
                        "question_type": question_type,
                        "topic": mcq.get("topic", "General"),
                        "difficulty": mcq.get("difficulty", "medium").lower(),
                        "bloom_level": mcq.get("bloom_level", "remember").lower(),
                        "source_doc": document_title,
                        "created_at": datetime.utcnow(),
                        "last_reviewed": None,
                        "times_attempted": 0,
                        "times_correct": 0
                    }
                    processed_mcqs.append(processed_mcq)
                
                return processed_mcqs[:question_count]  # Ensure exact count
        
        # Fallback if JSON parsing fails
        return create_fallback_mcqs(document_title, question_count, question_type, difficulty)
        
    except Exception as e:
        logger.error(f"Error generating MCQs: {e}")
        return create_fallback_mcqs(document_title, question_count, question_type, difficulty)

def create_fallback_mcqs(document_title: str, question_count: int, question_type: str = "single_correct", difficulty: str = "mixed") -> List[Dict[str, Any]]:
    """Create fallback MCQs if AI generation fails"""
    fallback_mcqs = []
    
    for i in range(min(question_count, 10)):  # Limit fallback to 10 questions max
        question_id = f"fallback_mcq_{i+1}_{uuid.uuid4().hex[:8]}"
        
        fallback_mcq = {
            "question_id": question_id,
            "question_text": f"What is the main topic of {document_title}?",
            "options": [
                {"option_text": f"This is a fallback question for {document_title}", "is_correct": True},
                {"option_text": "Please regenerate MCQs for better questions", "is_correct": False},
                {"option_text": "Review the document content", "is_correct": False},
                {"option_text": "Create custom questions manually", "is_correct": False}
            ],
            "explanation": f"This is a fallback MCQ for {document_title}. Please regenerate for better questions based on the actual content.",
            "question_type": question_type,
            "topic": "General",
            "difficulty": "easy",
            "bloom_level": "remember",
            "source_doc": document_title,
            "created_at": datetime.utcnow(),
            "last_reviewed": None,
            "times_attempted": 0,
            "times_correct": 0
        }
        fallback_mcqs.append(fallback_mcq)
    
    return fallback_mcqs

async def generate_podcast_script(file_path: str, file_type: str, document_title: str, podcast_length: str = "medium", content_focus: str = "full_document") -> str:
    """Generate podcast script using OpenAI"""
    try:
        # Read document content
        document_text = await extract_text_from_file(file_path, file_type)
        
        # Determine script length and style based on podcast_length
        length_instruction = ""
        if podcast_length == "quick":
            length_instruction = "Create a 2-3 minute podcast script focusing on key takeaways and essential points."
        elif podcast_length == "comprehensive":
            length_instruction = "Create a 7-10 minute podcast script with detailed explanations and comprehensive coverage."
        else:  # medium
            length_instruction = "Create a 4-6 minute podcast script balancing key concepts with engaging explanations."
            
        focus_instruction = ""
        if content_focus == "key_concepts":
            focus_instruction = "Focus primarily on the most important concepts, definitions, and core ideas."
        elif content_focus == "summary":
            focus_instruction = "Provide a comprehensive summary hitting all major points concisely."
        else:  # full_document
            focus_instruction = "Cover the full document content in an engaging, structured manner."
        
        podcast_prompt = f"""
        Please analyze this document and create an engaging podcast script based on the content.
        
        {length_instruction}
        {focus_instruction}
        
        Return the response as a JSON object with the following structure:
        
        {{
            "script": {{
                "title": "Engaging podcast title based on the content",
                "description": "Brief description of what the podcast covers",
                "intro": "Welcome introduction that hooks the listener",
                "main_content": "Main podcast content in conversational, natural speech",
                "conclusion": "Summary and closing remarks",
                "estimated_duration": "estimated duration in MM:SS format"
            }}
        }}
        
        SCRIPT WRITING GUIDELINES:
        1. Write in a conversational, engaging tone as if speaking directly to the listener
        2. Use natural speech patterns with appropriate pauses (indicated by commas and periods)
        3. Include smooth transitions between topics
        4. Explain complex concepts in simple, accessible language
        5. Use rhetorical questions to engage the listener
        6. Include brief examples or analogies when helpful
        7. Maintain an educational yet entertaining style
        8. Structure the content logically with clear flow
        9. Include verbal signposts like "First," "Next," "Finally," etc.
        10. End with actionable takeaways or thought-provoking questions
        
        INTRO TEMPLATE:
        "Welcome to this learning session on [topic]. I'm excited to share some fascinating insights about [subject] that will help you understand [key benefit]. Let's dive right in!"
        
        CONCLUSION TEMPLATE:
        "To wrap up today's session, we've covered [key points]. The main takeaway is [core message]. I hope this helps you [practical application]. Thanks for listening, and keep learning!"
        
        Document title: {document_title}
        Document content: {document_text[:4000]}...
        
        Create a script that transforms this written content into an engaging audio learning experience.
        """
        
        # Get AI response from OpenAI
        response = await openai_client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {
                    "role": "system",
                    "content": "You are an expert podcast script writer specializing in educational content. Your task is to transform written documents into engaging, conversational audio scripts that are perfect for listening. Focus on creating natural speech patterns, clear explanations, and an engaging narrative flow."
                },
                {
                    "role": "user",
                    "content": podcast_prompt
                }
            ],
            max_tokens=8000,
            temperature=0.7
        )
        
        ai_response = response.choices[0].message.content
        
        # Clean and parse JSON response
        json_start = ai_response.find('{')
        json_end = ai_response.rfind('}') + 1
        
        if json_start != -1 and json_end != 0:
            json_content = ai_response[json_start:json_end]
            script_data = json.loads(json_content)
            
            if "script" in script_data:
                script = script_data["script"]
                
                # Combine all script parts
                full_script = f"""
{script.get('intro', 'Welcome to this learning session.')}

{script.get('main_content', 'Here is the main content of our discussion.')}

{script.get('conclusion', 'Thank you for listening. Keep learning!')}
"""
                return {
                    "title": script.get("title", document_title + " - Podcast"),
                    "description": script.get("description", f"Educational podcast based on {document_title}"),
                    "script_text": full_script.strip(),
                    "estimated_duration": script.get("estimated_duration", "05:00")
                }
        
        # Fallback if JSON parsing fails
        return create_fallback_podcast_script(document_title, podcast_length)
        
    except Exception as e:
        logger.error(f"Error generating podcast script: {e}")
        return create_fallback_podcast_script(document_title, podcast_length)

def create_fallback_podcast_script(document_title: str, podcast_length: str = "medium") -> Dict[str, str]:
    """Create fallback podcast script if AI generation fails"""
    duration_map = {"quick": "03:00", "medium": "05:00", "comprehensive": "08:00"}
    
    script = f"""
Welcome to this learning session on {document_title}. 

Unfortunately, we couldn't generate a custom podcast script from your document content at this time. 
This is a fallback script to ensure you still have a functional podcast.

To get the best results, please try regenerating the podcast, or check that your document contains 
sufficient text content for script generation.

The main topics from {document_title} would typically be covered here, with engaging explanations 
and practical insights to help you learn effectively.

Thank you for using our podcast generation feature. We're constantly improving to bring you 
better educational content. Keep learning!
"""
    
    return {
        "title": f"{document_title} - Learning Podcast",
        "description": f"Educational podcast generated from {document_title}",
        "script_text": script.strip(),
        "estimated_duration": duration_map.get(podcast_length, "05:00")
    }

async def generate_audio_from_script(script_text: str, voice_style: str = "neutral", voice_gender: str = "female", voice_accent: str = "american", output_path: str = None) -> Dict[str, Any]:
    """Generate audio file from script using edge-tts with accent support"""
    try:
        import edge_tts
        import asyncio
        from pydub import AudioSegment
        import os
        
        # Extended voice map with accents
        voice_map = {
            # American English
            ("american", "female", "neutral"): "en-US-AriaNeural",
            ("american", "female", "enthusiastic"): "en-US-JennyNeural", 
            ("american", "female", "formal"): "en-US-SaraNeural",
            ("american", "female", "conversational"): "en-US-AriaNeural",
            ("american", "male", "neutral"): "en-US-GuyNeural",
            ("american", "male", "enthusiastic"): "en-US-BrianNeural",
            ("american", "male", "formal"): "en-US-DavisNeural", 
            ("american", "male", "conversational"): "en-US-GuyNeural",
            
            # British English
            ("british", "female", "neutral"): "en-GB-SoniaNeural",
            ("british", "female", "enthusiastic"): "en-GB-LibbyNeural", 
            ("british", "female", "formal"): "en-GB-SoniaNeural",
            ("british", "female", "conversational"): "en-GB-MaisieNeural",
            ("british", "male", "neutral"): "en-GB-RyanNeural",
            ("british", "male", "enthusiastic"): "en-GB-ThomasNeural",
            ("british", "male", "formal"): "en-GB-RyanNeural", 
            ("british", "male", "conversational"): "en-GB-AlfieNeural",
            
            # Indian English
            ("indian", "female", "neutral"): "en-IN-NeerjaNeural",
            ("indian", "female", "enthusiastic"): "en-IN-NeerjaNeural", 
            ("indian", "female", "formal"): "en-IN-NeerjaNeural",
            ("indian", "female", "conversational"): "en-IN-NeerjaNeural",
            ("indian", "male", "neutral"): "en-IN-PrabhatNeural",
            ("indian", "male", "enthusiastic"): "en-IN-PrabhatNeural",
            ("indian", "male", "formal"): "en-IN-PrabhatNeural", 
            ("indian", "male", "conversational"): "en-IN-PrabhatNeural",
            
            # Australian English
            ("australian", "female", "neutral"): "en-AU-NatashaNeural",
            ("australian", "female", "enthusiastic"): "en-AU-NatashaNeural", 
            ("australian", "female", "formal"): "en-AU-NatashaNeural",
            ("australian", "female", "conversational"): "en-AU-NatashaNeural",
            ("australian", "male", "neutral"): "en-AU-WilliamNeural",
            ("australian", "male", "enthusiastic"): "en-AU-WilliamNeural",
            ("australian", "male", "formal"): "en-AU-WilliamNeural", 
            ("australian", "male", "conversational"): "en-AU-WilliamNeural",
            
            # Canadian English
            ("canadian", "female", "neutral"): "en-CA-ClaraNeural",
            ("canadian", "female", "enthusiastic"): "en-CA-ClaraNeural", 
            ("canadian", "female", "formal"): "en-CA-ClaraNeural",
            ("canadian", "female", "conversational"): "en-CA-ClaraNeural",
            ("canadian", "male", "neutral"): "en-CA-LiamNeural",
            ("canadian", "male", "enthusiastic"): "en-CA-LiamNeural",
            ("canadian", "male", "formal"): "en-CA-LiamNeural", 
            ("canadian", "male", "conversational"): "en-CA-LiamNeural"
        }
        
        # Get appropriate voice
        voice_key = (voice_accent.lower(), voice_gender.lower(), voice_style.lower())
        voice_name = voice_map.get(voice_key, "en-US-AriaNeural")  # Default fallback
        
        # Generate unique filename if not provided
        if not output_path:
            audio_filename = f"podcast_{uuid.uuid4().hex[:8]}.mp3"
            output_path = os.path.join(EXPORT_DIR, audio_filename)
        
        # Ensure export directory exists
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        
        # Generate audio using edge-tts
        communicate = edge_tts.Communicate(script_text, voice_name)
        await communicate.save(output_path)
        
        # Get audio info using pydub
        audio = AudioSegment.from_mp3(output_path)
        duration_seconds = len(audio) / 1000
        file_size = os.path.getsize(output_path)
        
        # Format duration
        minutes = int(duration_seconds // 60)
        seconds = int(duration_seconds % 60)
        duration_formatted = f"{minutes:02d}:{seconds:02d}"
        
        return {
            "audio_file_path": output_path,
            "duration_seconds": int(duration_seconds),
            "duration_formatted": duration_formatted,
            "file_size_bytes": file_size,
            "voice_used": voice_name
        }
        
    except Exception as e:
        logger.error(f"Error generating audio: {e}")
        raise Exception(f"Audio generation failed: {str(e)}")

async def convert_audio_format(input_path: str, output_format: str = "wav") -> str:
    """Convert audio file to different format using pydub"""
    try:
        logger.info(f"Converting audio from {input_path} to {output_format}")
        
        from pydub import AudioSegment
        import os
        
        # Verify input file exists
        if not os.path.exists(input_path):
            raise FileNotFoundError(f"Input audio file not found: {input_path}")
        
        # Load the audio file
        logger.info(f"Loading audio file: {input_path}")
        audio = AudioSegment.from_mp3(input_path)
        
        # Generate output path
        base_path = os.path.splitext(input_path)[0]
        output_path = f"{base_path}.{output_format}"
        
        logger.info(f"Output path: {output_path}")
        
        # Export in the desired format
        if output_format.lower() == "wav":
            logger.info("Converting to WAV format")
            audio.export(output_path, format="wav")
        elif output_format.lower() == "mp3":
            logger.info("Converting to MP3 format")
            audio.export(output_path, format="mp3")
        else:
            raise ValueError(f"Unsupported format: {output_format}")
        
        # Verify output file was created
        if not os.path.exists(output_path):
            raise RuntimeError(f"Failed to create output file: {output_path}")
        
        logger.info(f"Successfully converted audio to: {output_path}")
        return output_path
        
    except Exception as e:
        logger.error(f"Error converting audio format: {str(e)}", exc_info=True)
        raise Exception(f"Audio conversion failed: {str(e)}")

def create_powerpoint_with_images(slides_data: List[Dict], theme: Dict, document_title: str) -> str:
    """Create PowerPoint presentation from slides data with embedded images"""
    try:
        logger.info(f"Starting PowerPoint creation for {len(slides_data)} slides with visual enhancements")
        
        # Create presentation
        prs = Presentation()
        
        # Set slide size (16:9 widescreen)
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # Parse colors safely
        def hex_to_rgb(hex_color: str):
            try:
                hex_color = hex_color.lstrip('#')
                if len(hex_color) != 6:
                    return (255, 255, 255)  # Default to white
                return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
            except ValueError:
                return (255, 255, 255)  # Default to white
        
        bg_color = hex_to_rgb(theme.get("background_color", "#FFFFFF"))
        default_title_color = hex_to_rgb(theme.get("title_color", "#000000"))
        default_text_color = hex_to_rgb(theme.get("text_color", "#333333"))
        
        for i, slide_data in enumerate(slides_data):
            logger.info(f"Creating slide {i+1}: {slide_data.get('title', 'Untitled')}")
            
            # Choose layout based on whether slide has image
            has_image = slide_data.get("image_id") is not None
            if has_image:
                slide_layout = prs.slide_layouts[8] if len(prs.slide_layouts) > 8 else prs.slide_layouts[1]  # Two Content layout for image + text
            else:
                slide_layout = prs.slide_layouts[1]  # Title and Content layout
            
            slide = prs.slides.add_slide(slide_layout)
            
            # Set background color
            try:
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(*bg_color)
            except Exception as e:
                logger.warning(f"Could not set background color: {e}")
            
            # Title with custom formatting
            if slide.shapes.title:
                title_shape = slide.shapes.title
                title_shape.text = slide_data.get("title", f"Slide {i+1}")
                
                try:
                    title_paragraph = title_shape.text_frame.paragraphs[0]
                    title_formatting = slide_data.get("title_formatting", {})
                    
                    title_paragraph.font.size = Pt(title_formatting.get("font_size", 28))
                    title_color = hex_to_rgb(title_formatting.get("text_color", theme.get("title_color", "#000000")))
                    title_paragraph.font.color.rgb = RGBColor(*title_color)
                    title_paragraph.font.name = title_formatting.get("font_family", theme.get("font_family", "Calibri"))
                    title_paragraph.font.bold = title_formatting.get("is_bold", True)
                    title_paragraph.font.italic = title_formatting.get("is_italic", False)
                except Exception as e:
                    logger.warning(f"Could not format title: {e}")
            
            # Content with custom formatting (adjusted for image layout)
            content_placeholder = None
            for shape in slide.shapes.placeholders:
                if shape.placeholder_format.idx == 1:  # Content placeholder
                    content_placeholder = shape
                    break
            
            if content_placeholder:
                try:
                    content_frame = content_placeholder.text_frame
                    content_frame.clear()
                    
                    content_points = slide_data.get("content", [])
                    body_formatting = slide_data.get("body_formatting", {})
                    
                    for j, point in enumerate(content_points):
                        if j == 0:
                            p = content_frame.paragraphs[0]
                        else:
                            p = content_frame.add_paragraph()
                        
                        p.text = f"• {str(point)}"
                        p.font.size = Pt(body_formatting.get("font_size", theme.get("font_size", 18)))
                        body_color = hex_to_rgb(body_formatting.get("text_color", theme.get("text_color", "#333333")))
                        p.font.color.rgb = RGBColor(*body_color)
                        p.font.name = body_formatting.get("font_family", theme.get("font_family", "Calibri"))
                        p.font.bold = body_formatting.get("is_bold", False)
                        p.font.italic = body_formatting.get("is_italic", False)
                        p.level = 0
                    
                    # Adjust content placeholder size if image is present
                    if has_image:
                        content_placeholder.width = Inches(6)  # Make room for image
                            
                except Exception as e:
                    logger.warning(f"Could not format content: {e}")
            
            # Add speaker notes
            try:
                notes_slide = slide.notes_slide
                notes_text_frame = notes_slide.notes_text_frame
                notes_text_frame.text = slide_data.get("speaker_notes", "")
            except Exception as e:
                logger.warning(f"Could not add speaker notes: {e}")
        
        # Save presentation with unique filename
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        filename = f"{document_title}_visual_{timestamp}.pptx"
        # Clean filename to avoid issues
        filename = "".join(c for c in filename if c.isalnum() or c in (' ', '-', '_')).rstrip()
        filename = filename.replace(' ', '_') + '.pptx'
        
        filepath = os.path.join(EXPORT_DIR, filename)
        
        logger.info(f"Saving enhanced PowerPoint to: {filepath}")
        prs.save(filepath)
        
        # Verify file was created
        if os.path.exists(filepath):
            file_size = os.path.getsize(filepath)
            logger.info(f"Enhanced PowerPoint created successfully. Size: {file_size} bytes")
            return filepath
        else:
            raise Exception("PowerPoint file was not created")
    
    except Exception as e:
        logger.error(f"Error creating enhanced PowerPoint: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to create PowerPoint: {str(e)}")

async def create_powerpoint_with_images_async(slides_data: List[Dict], theme: Dict, document_title: str) -> str:
    """Async wrapper for PowerPoint creation with image handling"""
    try:
        # Handle image retrieval asynchronously
        processed_slides = []
        for slide_data in slides_data:
            processed_slide = slide_data.copy()
            
            # If slide has image, get image data
            if slide_data.get("image_id"):
                try:
                    image_doc = await slide_images_collection.find_one({"image_id": slide_data["image_id"]})
                    if image_doc:
                        # Decode base64 image and save temporarily
                        image_data = base64.b64decode(image_doc["image_base64"])
                        temp_image_path = os.path.join(IMAGES_DIR, f"temp_{slide_data['image_id']}.png")
                        
                        with open(temp_image_path, "wb") as f:
                            f.write(image_data)
                        
                        processed_slide["temp_image_path"] = temp_image_path
                        logger.info(f"Prepared image for slide {slide_data.get('slide_number', 'unknown')}")
                        
                except Exception as e:
                    logger.warning(f"Could not prepare image for slide: {e}")
                    processed_slide["temp_image_path"] = None
            
            processed_slides.append(processed_slide)
        
        # Create PowerPoint synchronously with prepared data
        return create_powerpoint_with_images_sync(processed_slides, theme, document_title)
        
    except Exception as e:
        logger.error(f"Error in async PowerPoint creation: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to create PowerPoint: {str(e)}")

def create_powerpoint_with_images_sync(slides_data: List[Dict], theme: Dict, document_title: str) -> str:
    """Create PowerPoint presentation from slides data with prepared images"""
    try:
        logger.info(f"Starting PowerPoint creation for {len(slides_data)} slides with visual enhancements")
        
        # Create presentation
        prs = Presentation()
        
        # Set slide size (16:9 widescreen)
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # Parse colors safely
        def hex_to_rgb(hex_color: str):
            try:
                hex_color = hex_color.lstrip('#')
                if len(hex_color) != 6:
                    return (255, 255, 255)
                return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
            except ValueError:
                return (255, 255, 255)
        
        bg_color = hex_to_rgb(theme.get("background_color", "#FFFFFF"))
        
        for i, slide_data in enumerate(slides_data):
            logger.info(f"Creating slide {i+1}: {slide_data.get('title', 'Untitled')}")
            
            # Choose layout based on whether slide has image
            has_image = slide_data.get("temp_image_path") is not None
            if has_image:
                # Use a layout that supports images and text
                slide_layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[1]  # Blank or content layout
            else:
                slide_layout = prs.slide_layouts[1]  # Title and Content layout
            
            slide = prs.slides.add_slide(slide_layout)
            
            # Set background color
            try:
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(*bg_color)
            except Exception as e:
                logger.warning(f"Could not set background color: {e}")
            
            # Add title manually for better control
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
            title_frame = title_box.text_frame
            title_frame.text = slide_data.get("title", f"Slide {i+1}")
            
            # Format title
            try:
                title_paragraph = title_frame.paragraphs[0]
                title_formatting = slide_data.get("title_formatting", {})
                
                title_paragraph.font.size = Pt(title_formatting.get("font_size", 28))
                title_color = hex_to_rgb(title_formatting.get("text_color", "#000000"))
                title_paragraph.font.color.rgb = RGBColor(*title_color)
                title_paragraph.font.name = title_formatting.get("font_family", "Calibri")
                title_paragraph.font.bold = title_formatting.get("is_bold", True)
                title_paragraph.font.italic = title_formatting.get("is_italic", False)
            except Exception as e:
                logger.warning(f"Could not format title: {e}")
            
            # Add image if available
            if has_image and slide_data.get("temp_image_path"):
                try:
                    if os.path.exists(slide_data["temp_image_path"]):
                        # Add image to right side
                        left = Inches(7)
                        top = Inches(2)
                        width = Inches(5.5)
                        height = Inches(4.5)
                        
                        slide.shapes.add_picture(slide_data["temp_image_path"], left, top, width, height)
                        logger.info(f"Added image to slide {i+1}")
                        
                        # Clean up temp file
                        os.remove(slide_data["temp_image_path"])
                        
                except Exception as e:
                    logger.warning(f"Could not add image to slide {i+1}: {e}")
            
            # Add content text
            try:
                content_left = Inches(0.5)
                content_top = Inches(2)
                content_width = Inches(6) if has_image else Inches(12)
                content_height = Inches(4.5)
                
                content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
                content_frame = content_box.text_frame
                
                content_points = slide_data.get("content", [])
                body_formatting = slide_data.get("body_formatting", {})
                
                for j, point in enumerate(content_points):
                    if j == 0:
                        p = content_frame.paragraphs[0]
                    else:
                        p = content_frame.add_paragraph()
                    
                    p.text = f"• {str(point)}"
                    p.font.size = Pt(body_formatting.get("font_size", 18))
                    body_color = hex_to_rgb(body_formatting.get("text_color", "#333333"))
                    p.font.color.rgb = RGBColor(*body_color)
                    p.font.name = body_formatting.get("font_family", "Calibri")
                    p.font.bold = body_formatting.get("is_bold", False)
                    p.font.italic = body_formatting.get("is_italic", False)
                    p.level = 0
                        
            except Exception as e:
                logger.warning(f"Could not format content: {e}")
            
            # Add speaker notes
            try:
                notes_slide = slide.notes_slide
                notes_text_frame = notes_slide.notes_text_frame
                notes_text_frame.text = slide_data.get("speaker_notes", "")
            except Exception as e:
                logger.warning(f"Could not add speaker notes: {e}")
        
        # Save presentation
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        filename = f"{document_title}_visual_{timestamp}.pptx"
        filename = "".join(c for c in filename if c.isalnum() or c in (' ', '-', '_')).rstrip()
        filename = filename.replace(' ', '_') + '.pptx'
        
        filepath = os.path.join(EXPORT_DIR, filename)
        
        logger.info(f"Saving enhanced PowerPoint to: {filepath}")
        prs.save(filepath)
        
        if os.path.exists(filepath):
            file_size = os.path.getsize(filepath)
            logger.info(f"Enhanced PowerPoint created successfully. Size: {file_size} bytes")
            return filepath
        else:
            raise Exception("PowerPoint file was not created")
    
    except Exception as e:
        logger.error(f"Error creating enhanced PowerPoint: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to create PowerPoint: {str(e)}")

# API Routes
@app.get("/api/health")
async def health_check():
    return {"status": "healthy", "service": "EduMind AI Backend with Visual Enhancement"}

@app.post("/api/upload")
async def upload_document(file: UploadFile = File(...)):
    """Upload and process document"""
    try:
        # Validate file type
        allowed_types = ["application/pdf", "application/vnd.openxmlformats-officedocument.wordprocessingml.document", "text/plain"]
        if file.content_type not in allowed_types:
            raise HTTPException(status_code=400, detail="File type not supported. Please upload PDF, DOCX, or TXT files.")
        
        # Generate unique file ID
        file_id = str(uuid.uuid4())
        file_extension = file.filename.split('.')[-1].lower()
        file_path = os.path.join(UPLOAD_DIR, f"{file_id}.{file_extension}")
        
        # Save file
        async with aiofiles.open(file_path, 'wb') as f:
            content = await file.read()
            await f.write(content)
        
        # Store document info in database
        doc_info = {
            "file_id": file_id,
            "filename": file.filename,
            "file_path": file_path,
            "file_type": file.content_type,
            "uploaded_at": datetime.utcnow()
        }
        await documents_collection.insert_one(doc_info)
        
        return {
            "file_id": file_id,
            "filename": file.filename,
            "message": "File uploaded successfully"
        }
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Upload failed: {str(e)}")

@app.post("/api/generate-slides")
async def generate_slides(request: SlideGenerationRequest):
    """Generate presentation slides from uploaded document with optional visual enhancements"""
    try:
        # Validate slide count
        if request.slide_count < 3 or request.slide_count > 20:
            raise HTTPException(status_code=400, detail="Slide count must be between 3 and 20")
        
        # Get document info
        doc = await documents_collection.find_one({"file_id": request.file_id})
        if not doc:
            raise HTTPException(status_code=404, detail="Document not found")
        
        file_path = doc["file_path"]
        file_type = doc["file_type"]
        filename = doc["filename"]
        
        # Check if slides already exist for this document
        existing_slides = await slides_collection.find_one({"file_id": request.file_id})
        if existing_slides:
            # Delete existing slides to regenerate with new settings
            await slides_collection.delete_one({"file_id": request.file_id})
            # Also clean up associated images
            await slide_images_collection.delete_many({"slide_id": {"$regex": f"^slide_.*"}})
        
        # Generate slides using AI with specified count and visual options
        slides_data = await generate_slide_content(
            file_path, 
            file_type, 
            filename, 
            request.slide_count,
            request.include_images
        )
        
        # Prepare slide presentation data
        slide_presentation = {
            "file_id": request.file_id,
            "session_id": request.session_id,
            "slides": slides_data,
            "total_slides": len(slides_data),
            "document_title": filename.rsplit('.', 1)[0],
            "generated_at": datetime.utcnow(),
            "is_editable": True,
            "has_images": request.include_images
        }
        
        # Save slides to database
        await slides_collection.insert_one(slide_presentation)
        
        return {
            "slides": slides_data,
            "total_slides": len(slides_data),
            "document_title": filename.rsplit('.', 1)[0],
            "generated_at": datetime.utcnow(),
            "is_editable": True,
            "has_images": request.include_images
        }
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Slide generation failed: {str(e)}")

@app.post("/api/generate-slide-image")
async def generate_slide_image_endpoint(request: ImageGenerationRequest):
    """Generate or regenerate an image for a specific slide"""
    try:
        # Find slide data
        slide_doc = await slides_collection.find_one({"slides.slide_id": request.slide_id})
        if not slide_doc:
            raise HTTPException(status_code=404, detail="Slide not found")
        
        # Get slide info
        slide_info = None
        for slide in slide_doc["slides"]:
            if slide["slide_id"] == request.slide_id:
                slide_info = slide
                break
        
        if not slide_info:
            raise HTTPException(status_code=404, detail="Slide not found")
        
        # Generate image
        image_base64 = await generate_slide_image(
            slide_info["title"], 
            slide_info["content"], 
            request.prompt
        )
        
        if not image_base64:
            raise HTTPException(status_code=500, detail="Failed to generate image")
        
        # Store image in database
        image_id = f"img_{request.slide_id}_{uuid.uuid4().hex[:8]}"
        image_data = {
            "image_id": image_id,
            "slide_id": request.slide_id,
            "image_base64": image_base64,
            "prompt_used": request.prompt or f"Educational illustration for '{slide_info['title']}'",
            "generated_at": datetime.utcnow()
        }
        await slide_images_collection.insert_one(image_data)
        
        # Update slide with image reference
        slides = slide_doc["slides"]
        for slide in slides:
            if slide["slide_id"] == request.slide_id:
                slide["image_id"] = image_id
                break
        
        await slides_collection.update_one(
            {"_id": slide_doc["_id"]},
            {"$set": {"slides": slides, "updated_at": datetime.utcnow()}}
        )
        
        return {
            "image_id": image_id,
            "image_base64": image_base64,
            "message": "Image generated successfully"
        }
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Image generation failed: {str(e)}")

@app.get("/api/slide-image/{image_id}")
async def get_slide_image(image_id: str):
    """Retrieve a slide image by ID"""
    try:
        image_doc = await slide_images_collection.find_one({"image_id": image_id})
        if not image_doc:
            raise HTTPException(status_code=404, detail="Image not found")
        
        return {
            "image_id": image_doc["image_id"],
            "image_base64": image_doc["image_base64"],
            "prompt_used": image_doc["prompt_used"],
            "generated_at": image_doc["generated_at"]
        }
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to retrieve image: {str(e)}")

@app.put("/api/slides/edit")
async def edit_slide(request: SlideEditRequest):
    """Edit individual slide content with separate title and body formatting"""
    try:
        # Find the slide in the database
        slide_doc = await slides_collection.find_one({"slides.slide_id": request.slide_id})
        if not slide_doc:
            raise HTTPException(status_code=404, detail="Slide not found")
        
        # Update the specific slide
        slides = slide_doc["slides"]
        for slide in slides:
            if slide["slide_id"] == request.slide_id:
                slide["title"] = request.title
                slide["content"] = request.content
                slide["speaker_notes"] = request.speaker_notes
                if request.title_formatting:
                    slide["title_formatting"] = request.title_formatting.dict()
                if request.body_formatting:
                    slide["body_formatting"] = request.body_formatting.dict()
                break
        
        # Update in database
        await slides_collection.update_one(
            {"_id": slide_doc["_id"]},
            {"$set": {"slides": slides, "updated_at": datetime.utcnow()}}
        )
        
        return {"message": "Slide updated successfully"}
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Slide update failed: {str(e)}")

@app.get("/api/slides/{file_id}")
async def get_slides(file_id: str):
    """Get generated slides for a document"""
    try:
        slides_doc = await slides_collection.find_one({"file_id": file_id})
        if not slides_doc:
            raise HTTPException(status_code=404, detail="Slides not found for this document")
        
        return {
            "slides": slides_doc["slides"],
            "total_slides": slides_doc["total_slides"],
            "document_title": slides_doc["document_title"],
            "generated_at": slides_doc["generated_at"],
            "is_editable": slides_doc.get("is_editable", True),
            "has_images": slides_doc.get("has_images", False)
        }
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to retrieve slides: {str(e)}")

@app.get("/api/presentation-themes")
async def get_presentation_themes():
    """Get available presentation themes with visual design elements"""
    return {"themes": PRESENTATION_THEMES}

@app.post("/api/export-presentation")
async def export_presentation(request: ExportRequest):
    """Export slides as PowerPoint presentation with images"""
    try:
        logger.info(f"Export request received for file_id: {request.file_id}")
        
        # Get slides data
        slides_doc = await slides_collection.find_one({"file_id": request.file_id})
        if not slides_doc:
            raise HTTPException(status_code=404, detail="Slides not found")
        
        slides_data = slides_doc["slides"]
        document_title = slides_doc["document_title"]
        
        logger.info(f"Creating enhanced PowerPoint for {len(slides_data)} slides")
        
        # Create PowerPoint file with images
        ppt_path = await create_powerpoint_with_images_async(slides_data, request.theme.dict(), document_title)
        
        if not os.path.exists(ppt_path):
            raise HTTPException(status_code=500, detail="PowerPoint file was not created")
        
        # Return file for download
        filename = os.path.basename(ppt_path)
        logger.info(f"Sending enhanced file: {filename}")
        
        return FileResponse(
            path=ppt_path,
            filename=filename,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f"attachment; filename={filename}"}
        )
    
    except Exception as e:
        logger.error(f"Enhanced export failed: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Export failed: {str(e)}")

# Continue with existing chat and other endpoints...
@app.post("/api/chat")
async def chat_with_ai(chat_data: ChatMessage):
    """Chat with AI about uploaded documents"""
    try:
        # Get or create chat session
        session = await get_or_create_chat_session(chat_data.session_id)
        
        # Prepare user message
        user_message_text = chat_data.message
        document_context = ""
        
        # If file_id is provided, get document content for context
        if chat_data.file_id:
            doc = await documents_collection.find_one({"file_id": chat_data.file_id})
            if doc:
                file_path = doc["file_path"]
                file_type = doc["file_type"]
                
                # Extract document text for context
                document_text = await extract_text_from_file(file_path, file_type)
                document_context = f"\n\nDocument context (from {doc['filename']}):\n{document_text[:2000]}..."
        
        # Create enhanced prompt with document context
        enhanced_message = f"{user_message_text}{document_context}"
        
        # Get AI response from OpenAI
        response = await openai_client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {
                    "role": "system",
                    "content": "You are EduMind AI, an intelligent learning assistant. You help users understand and learn from their documents by providing clear, educational responses. Always be helpful, accurate, and focus on educational value. When document context is provided, use it to give specific, relevant answers."
                },
                {
                    "role": "user",
                    "content": enhanced_message
                }
            ],
            max_tokens=4096,
            temperature=0.7
        )
        
        ai_response = response.choices[0].message.content
        
        # Save user message to database
        user_msg_data = {
            "session_id": chat_data.session_id,
            "message": user_message_text,
            "sender": "user",
            "timestamp": datetime.utcnow(),
            "file_id": chat_data.file_id
        }
        await messages_collection.insert_one(user_msg_data)
        
        # Save AI response to database
        ai_msg_data = {
            "session_id": chat_data.session_id,
            "message": ai_response,
            "sender": "ai",
            "timestamp": datetime.utcnow()
        }
        await messages_collection.insert_one(ai_msg_data)
        
        # Update session
        await chat_sessions_collection.update_one(
            {"session_id": chat_data.session_id},
            {"$set": {"updated_at": datetime.utcnow()}}
        )
        
        return {
            "response": ai_response,
            "session_id": chat_data.session_id,
            "timestamp": datetime.utcnow()
        }
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Chat failed: {str(e)}")

@app.get("/api/sessions")
async def get_chat_sessions():
    """Get all chat sessions"""
    sessions = await chat_sessions_collection.find().sort("updated_at", -1).limit(20).to_list(length=20)
    for session in sessions:
        session["_id"] = str(session["_id"])
    return {"sessions": sessions}

@app.get("/api/sessions/{session_id}/messages")
async def get_session_messages(session_id: str):
    """Get messages for a specific session"""
    messages = await messages_collection.find({"session_id": session_id}).sort("timestamp", 1).to_list(length=100)
    for message in messages:
        message["_id"] = str(message["_id"])
    return {"messages": messages}

@app.get("/api/documents")
async def get_uploaded_documents():
    """Get list of uploaded documents"""
    docs = await documents_collection.find().sort("uploaded_at", -1).limit(20).to_list(length=20)
    for doc in docs:
        doc["_id"] = str(doc["_id"])
    return {"documents": docs}

# Flashcard API Routes
@app.post("/api/generate-flashcards")
async def generate_flashcards(request: FlashcardGenerationRequest):
    """Generate flashcards from uploaded document using AI"""
    try:
        # Validate card count
        if request.card_count < 5 or request.card_count > 50:
            raise HTTPException(status_code=400, detail="Card count must be between 5 and 50")
        
        # Get document info
        doc = await documents_collection.find_one({"file_id": request.file_id})
        if not doc:
            raise HTTPException(status_code=404, detail="Document not found")
        
        file_path = doc["file_path"]
        file_type = doc["file_type"]
        filename = doc["filename"]
        
        # Check if flashcards already exist for this document
        existing_set = await flashcard_sets_collection.find_one({"file_id": request.file_id})
        if existing_set:
            # Delete existing flashcards to regenerate with new settings
            await flashcard_sets_collection.delete_one({"file_id": request.file_id})
            await flashcards_collection.delete_many({"file_id": request.file_id})
        
        # Generate flashcards using AI
        cards_data = await generate_flashcard_content(
            file_path, 
            file_type, 
            filename, 
            request.card_count,
            request.card_type
        )
        
        # Create flashcard set
        set_id = f"set_{uuid.uuid4()}"
        flashcard_set = {
            "set_id": set_id,
            "set_name": f"{filename.rsplit('.', 1)[0]} - Flashcards",
            "file_id": request.file_id,
            "session_id": request.session_id,
            "document_title": filename.rsplit('.', 1)[0],
            "total_cards": len(cards_data),
            "card_type": request.card_type,
            "generated_at": datetime.utcnow(),
            "last_modified": datetime.utcnow(),
            "is_editable": True
        }
        
        # Save flashcard set to database
        await flashcard_sets_collection.insert_one(flashcard_set)
        
        # Save individual flashcards with file_id reference
        for card in cards_data:
            card["file_id"] = request.file_id
        await flashcards_collection.insert_many(cards_data)
        
        # Remove MongoDB _id fields from response data
        clean_cards_data = []
        for card in cards_data:
            clean_card = {k: v for k, v in card.items() if k != "_id"}
            clean_cards_data.append(clean_card)
        
        return {
            "set_id": set_id,
            "flashcards": clean_cards_data,
            "total_cards": len(clean_cards_data),
            "set_name": flashcard_set["set_name"],
            "generated_at": datetime.utcnow(),
            "card_type": request.card_type
        }
    
    except Exception as e:
        logger.error(f"Flashcard generation error: {e}")
        raise HTTPException(status_code=500, detail=f"Flashcard generation failed: {str(e)}")

@app.get("/api/flashcards/{file_id}")
async def get_flashcards(file_id: str):
    """Get generated flashcards for a document"""
    try:
        # Get flashcard set
        flashcard_set = await flashcard_sets_collection.find_one({"file_id": file_id})
        if not flashcard_set:
            raise HTTPException(status_code=404, detail="Flashcards not found for this document")
        
        # Get flashcards
        flashcards = await flashcards_collection.find({"file_id": file_id}).sort("created_at", 1).to_list(length=100)
        
        # Remove MongoDB _id fields from flashcards
        clean_flashcards = []
        for card in flashcards:
            clean_card = {k: v for k, v in card.items() if k != "_id"}
            clean_flashcards.append(clean_card)
        
        return {
            "set_id": flashcard_set["set_id"],
            "set_name": flashcard_set["set_name"],
            "flashcards": clean_flashcards,
            "total_cards": flashcard_set["total_cards"],
            "card_type": flashcard_set["card_type"],
            "generated_at": flashcard_set["generated_at"],
            "last_modified": flashcard_set["last_modified"],
            "is_editable": flashcard_set.get("is_editable", True)
        }
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to retrieve flashcards: {str(e)}")

@app.put("/api/flashcards/edit")
async def edit_flashcard(request: FlashcardEditRequest):
    """Edit individual flashcard content"""
    try:
        # Find and update the flashcard
        result = await flashcards_collection.update_one(
            {"card_id": request.card_id},
            {"$set": {
                "question": request.question,
                "answer": request.answer,
                "topic": request.topic,
                "tags": request.tags,
                "difficulty": request.difficulty
            }}
        )
        
        if result.modified_count == 0:
            raise HTTPException(status_code=404, detail="Flashcard not found")
        
        # Update the set's last modified time
        card = await flashcards_collection.find_one({"card_id": request.card_id})
        if card:
            await flashcard_sets_collection.update_one(
                {"file_id": card["file_id"]},
                {"$set": {"last_modified": datetime.utcnow()}}
            )
        
        return {"message": "Flashcard updated successfully"}
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Flashcard update failed: {str(e)}")

@app.delete("/api/flashcards/{card_id}")
async def delete_flashcard(card_id: str):
    """Delete a specific flashcard"""
    try:
        # Get card info before deletion
        card = await flashcards_collection.find_one({"card_id": card_id})
        if not card:
            raise HTTPException(status_code=404, detail="Flashcard not found")
        
        # Delete the flashcard
        await flashcards_collection.delete_one({"card_id": card_id})
        
        # Update set total count and last modified
        remaining_count = await flashcards_collection.count_documents({"file_id": card["file_id"]})
        await flashcard_sets_collection.update_one(
            {"file_id": card["file_id"]},
            {"$set": {
                "total_cards": remaining_count,
                "last_modified": datetime.utcnow()
            }}
        )
        
        return {"message": "Flashcard deleted successfully"}
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Flashcard deletion failed: {str(e)}")

@app.post("/api/flashcards/export")
async def export_flashcards(request: FlashcardExportRequest):
    """Export flashcards in various formats"""
    try:
        # Get flashcard set
        flashcard_set = await flashcard_sets_collection.find_one({"set_id": request.set_id})
        if not flashcard_set:
            raise HTTPException(status_code=404, detail="Flashcard set not found")
        
        # Get flashcards
        flashcards = await flashcards_collection.find({"file_id": flashcard_set["file_id"]}).sort("created_at", 1).to_list(length=200)
        
        # Remove MongoDB _id fields from flashcards for export
        clean_flashcards = []
        for card in flashcards:
            clean_card = {k: v for k, v in card.items() if k != "_id"}
            clean_flashcards.append(clean_card)
        
        if request.export_format.lower() == "csv":
            return export_flashcards_csv(clean_flashcards, flashcard_set)
        elif request.export_format.lower() == "anki":
            return export_flashcards_anki(clean_flashcards, flashcard_set)
        elif request.export_format.lower() == "json":
            return export_flashcards_json(clean_flashcards, flashcard_set)
        else:
            raise HTTPException(status_code=400, detail="Unsupported export format. Use csv, anki, or json")
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Export failed: {str(e)}")

def export_flashcards_csv(flashcards: List[Dict], flashcard_set: Dict):
    """Export flashcards as CSV"""
    import csv
    from io import StringIO
    
    output = StringIO()
    writer = csv.writer(output)
    
    # Write header
    writer.writerow(['Question', 'Answer', 'Topic', 'Tags', 'Difficulty'])
    
    # Write flashcards
    for card in flashcards:
        writer.writerow([
            card.get("question", ""),
            card.get("answer", ""),
            card.get("topic", ""),
            ", ".join(card.get("tags", [])),
            card.get("difficulty", "")
        ])
    
    # Create response
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    filename = f"{flashcard_set['set_name']}_flashcards_{timestamp}.csv"
    
    from fastapi.responses import Response
    return Response(
        content=output.getvalue(),
        media_type="text/csv",
        headers={"Content-Disposition": f"attachment; filename={filename}"}
    )

def export_flashcards_anki(flashcards: List[Dict], flashcard_set: Dict):
    """Export flashcards in Anki-compatible format"""
    from io import StringIO
    
    output = StringIO()
    
    # Anki format: Question\tAnswer\tTags
    for card in flashcards:
        tags = " ".join(card.get("tags", []))
        if tags:
            tags = f"{card.get('topic', '')} {tags}".strip()
        else:
            tags = card.get('topic', '')
        
        output.write(f"{card.get('question', '')}\t{card.get('answer', '')}\t{tags}\n")
    
    # Create response
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    filename = f"{flashcard_set['set_name']}_anki_{timestamp}.txt"
    
    from fastapi.responses import Response
    return Response(
        content=output.getvalue(),
        media_type="text/plain",
        headers={"Content-Disposition": f"attachment; filename={filename}"}
    )

def export_flashcards_json(flashcards: List[Dict], flashcard_set: Dict):
    """Export flashcards as JSON"""
    # Convert datetime objects to ISO format strings for JSON serialization
    json_flashcards = []
    for card in flashcards:
        json_card = {}
        for key, value in card.items():
            if isinstance(value, datetime):
                json_card[key] = value.isoformat()
            else:
                json_card[key] = value
        json_flashcards.append(json_card)
    
    export_data = {
        "set_info": {
            "set_name": flashcard_set["set_name"],
            "total_cards": len(json_flashcards),
            "exported_at": datetime.utcnow().isoformat()
        },
        "flashcards": json_flashcards
    }
    
    # Create response
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    filename = f"{flashcard_set['set_name']}_flashcards_{timestamp}.json"
    
    return JSONResponse(
        content=export_data,
        headers={"Content-Disposition": f"attachment; filename={filename}"}
    )

# MCQ API Routes
@app.post("/api/generate-mcqs")
async def generate_mcqs(request: MCQGenerationRequest):
    """Generate MCQs from uploaded document using AI"""
    try:
        # Validate question count
        if request.question_count < 5 or request.question_count > 30:
            raise HTTPException(status_code=400, detail="Question count must be between 5 and 30")
        
        # Get document info
        doc = await documents_collection.find_one({"file_id": request.file_id})
        if not doc:
            raise HTTPException(status_code=404, detail="Document not found")
        
        file_path = doc["file_path"]
        file_type = doc["file_type"]
        filename = doc["filename"]
        
        # Check if MCQs already exist for this document
        existing_set = await mcq_sets_collection.find_one({"file_id": request.file_id})
        if existing_set:
            # Delete existing MCQs to regenerate with new settings
            await mcq_sets_collection.delete_one({"file_id": request.file_id})
            await mcqs_collection.delete_many({"file_id": request.file_id})
        
        # Generate MCQs using AI
        mcqs_data = await generate_mcq_content(
            file_path, 
            file_type, 
            filename, 
            request.question_count,
            request.question_type,
            request.difficulty
        )
        
        # Calculate difficulty distribution
        difficulty_dist = {}
        for mcq in mcqs_data:
            diff = mcq.get("difficulty", "medium")
            difficulty_dist[diff] = difficulty_dist.get(diff, 0) + 1
        
        # Create MCQ set
        set_id = f"mcq_set_{uuid.uuid4()}"
        mcq_set = {
            "set_id": set_id,
            "set_name": f"{filename.rsplit('.', 1)[0]} - MCQs",
            "file_id": request.file_id,
            "session_id": request.session_id,
            "document_title": filename.rsplit('.', 1)[0],
            "total_questions": len(mcqs_data),
            "question_type": request.question_type,
            "difficulty_distribution": difficulty_dist,
            "generated_at": datetime.utcnow(),
            "last_modified": datetime.utcnow(),
            "is_editable": True
        }
        
        # Save MCQ set to database
        await mcq_sets_collection.insert_one(mcq_set)
        
        # Save individual MCQs with file_id reference
        for mcq in mcqs_data:
            mcq["file_id"] = request.file_id
        await mcqs_collection.insert_many(mcqs_data)
        
        # Clean MCQs data for JSON response (remove MongoDB ObjectIds)
        clean_mcqs = []
        for mcq in mcqs_data:
            clean_mcq = mcq.copy()
            if "_id" in clean_mcq:
                del clean_mcq["_id"]
            # Convert datetime objects to strings
            if "created_at" in clean_mcq and clean_mcq["created_at"]:
                clean_mcq["created_at"] = clean_mcq["created_at"].isoformat()
            if "last_reviewed" in clean_mcq and clean_mcq["last_reviewed"]:
                clean_mcq["last_reviewed"] = clean_mcq["last_reviewed"].isoformat()
            clean_mcqs.append(clean_mcq)
        
        return {
            "set_id": set_id,
            "mcqs": clean_mcqs,
            "total_questions": len(clean_mcqs),
            "set_name": mcq_set["set_name"],
            "generated_at": datetime.utcnow().isoformat(),
            "question_type": request.question_type,
            "difficulty_distribution": difficulty_dist
        }
    
    except Exception as e:
        logger.error(f"MCQ generation error: {e}")
        raise HTTPException(status_code=500, detail=f"MCQ generation failed: {str(e)}")

@app.get("/api/mcqs/{file_id}")
async def get_mcqs(file_id: str):
    """Get generated MCQs for a document"""
    try:
        # Get MCQ set
        mcq_set = await mcq_sets_collection.find_one({"file_id": file_id})
        if not mcq_set:
            raise HTTPException(status_code=404, detail="MCQs not found for this document")
        
        # Get MCQs
        mcqs_cursor = await mcqs_collection.find({"file_id": file_id}).sort("created_at", 1).to_list(length=100)
        
        # Clean MCQs data for JSON response (remove MongoDB ObjectIds)
        clean_mcqs = []
        for mcq in mcqs_cursor:
            clean_mcq = mcq.copy()
            if "_id" in clean_mcq:
                del clean_mcq["_id"]
            # Convert datetime objects to strings
            if "created_at" in clean_mcq and clean_mcq["created_at"]:
                clean_mcq["created_at"] = clean_mcq["created_at"].isoformat()
            if "last_reviewed" in clean_mcq and clean_mcq["last_reviewed"]:
                clean_mcq["last_reviewed"] = clean_mcq["last_reviewed"].isoformat()
            clean_mcqs.append(clean_mcq)
        
        return {
            "set_id": mcq_set["set_id"],
            "set_name": mcq_set["set_name"],
            "mcqs": clean_mcqs,
            "total_questions": mcq_set["total_questions"],
            "question_type": mcq_set["question_type"],
            "difficulty_distribution": mcq_set.get("difficulty_distribution", {}),
            "generated_at": mcq_set["generated_at"].isoformat() if mcq_set["generated_at"] else None,
            "last_modified": mcq_set["last_modified"].isoformat() if mcq_set["last_modified"] else None,
            "is_editable": mcq_set.get("is_editable", True)
        }
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to retrieve MCQs: {str(e)}")

@app.put("/api/mcqs/edit")
async def edit_mcq(request: MCQEditRequest):
    """Edit individual MCQ content"""
    try:
        # Validate that exactly one option is correct
        correct_count = sum(1 for opt in request.options if opt.is_correct)
        if correct_count != 1:
            raise HTTPException(status_code=400, detail="Exactly one option must be marked as correct")
        
        # Find and update the MCQ
        result = await mcqs_collection.update_one(
            {"question_id": request.question_id},
            {"$set": {
                "question_text": request.question_text,
                "options": [opt.dict() for opt in request.options],
                "explanation": request.explanation,
                "topic": request.topic,
                "difficulty": request.difficulty,
                "bloom_level": request.bloom_level
            }}
        )
        
        if result.modified_count == 0:
            raise HTTPException(status_code=404, detail="MCQ not found")
        
        # Update the set's last modified time
        mcq = await mcqs_collection.find_one({"question_id": request.question_id})
        if mcq:
            await mcq_sets_collection.update_one(
                {"file_id": mcq["file_id"]},
                {"$set": {"last_modified": datetime.utcnow()}}
            )
        
        return {"message": "MCQ updated successfully"}
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"MCQ update failed: {str(e)}")

@app.delete("/api/mcqs/{question_id}")
async def delete_mcq(question_id: str):
    """Delete a specific MCQ"""
    try:
        # Get MCQ info before deletion
        mcq = await mcqs_collection.find_one({"question_id": question_id})
        if not mcq:
            raise HTTPException(status_code=404, detail="MCQ not found")
        
        # Delete the MCQ
        await mcqs_collection.delete_one({"question_id": question_id})
        
        # Update set total count and last modified
        remaining_count = await mcqs_collection.count_documents({"file_id": mcq["file_id"]})
        
        # Recalculate difficulty distribution
        remaining_mcqs = await mcqs_collection.find({"file_id": mcq["file_id"]}).to_list(length=100)
        difficulty_dist = {}
        for remaining_mcq in remaining_mcqs:
            diff = remaining_mcq.get("difficulty", "medium")
            difficulty_dist[diff] = difficulty_dist.get(diff, 0) + 1
        
        await mcq_sets_collection.update_one(
            {"file_id": mcq["file_id"]},
            {"$set": {
                "total_questions": remaining_count,
                "difficulty_distribution": difficulty_dist,
                "last_modified": datetime.utcnow()
            }}
        )
        
        return {"message": "MCQ deleted successfully"}
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"MCQ deletion failed: {str(e)}")

@app.post("/api/mcqs/export")
async def export_mcqs(request: MCQExportRequest):
    """Export MCQs in various formats"""
    try:
        # Get MCQ set
        mcq_set = await mcq_sets_collection.find_one({"set_id": request.set_id})
        if not mcq_set:
            raise HTTPException(status_code=404, detail="MCQ set not found")
        
        # Get MCQs
        mcqs = await mcqs_collection.find({"file_id": mcq_set["file_id"]}).sort("created_at", 1).to_list(length=200)
        
        if request.export_format.lower() == "csv":
            return export_mcqs_csv(mcqs, mcq_set)
        elif request.export_format.lower() == "json":
            return export_mcqs_json(mcqs, mcq_set)
        elif request.export_format.lower() == "pdf":
            return export_mcqs_pdf(mcqs, mcq_set)
        else:
            raise HTTPException(status_code=400, detail="Unsupported export format. Use csv, pdf, or json")
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Export failed: {str(e)}")

def export_mcqs_csv(mcqs: List[Dict], mcq_set: Dict):
    """Export MCQs as CSV"""
    import csv
    from io import StringIO
    
    output = StringIO()
    writer = csv.writer(output)
    
    # Write header
    writer.writerow(['Question', 'Option A', 'Option B', 'Option C', 'Option D', 'Correct Answer', 'Explanation', 'Topic', 'Difficulty', 'Bloom Level'])
    
    # Write MCQs
    for mcq in mcqs:
        options = mcq.get("options", [])
        correct_answer = ""
        
        # Find correct answer
        for i, opt in enumerate(options):
            if opt.get("is_correct", False):
                correct_answer = chr(65 + i)  # A, B, C, D
                break
        
        # Ensure we have 4 options
        option_texts = ["", "", "", ""]
        for i, opt in enumerate(options[:4]):
            option_texts[i] = opt.get("option_text", "")
        
        writer.writerow([
            mcq.get("question_text", ""),
            option_texts[0],
            option_texts[1], 
            option_texts[2],
            option_texts[3],
            correct_answer,
            mcq.get("explanation", ""),
            mcq.get("topic", ""),
            mcq.get("difficulty", ""),
            mcq.get("bloom_level", "")
        ])
    
    # Create response
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    # Clean set name for filename
    clean_set_name = mcq_set.get("set_name", "mcqs").replace(" ", "_")
    filename = f"{clean_set_name}_{timestamp}.csv"
    
    from fastapi.responses import Response
    return Response(
        content=output.getvalue(),
        media_type="text/csv",
        headers={"Content-Disposition": f"attachment; filename={filename}"}
    )

def export_mcqs_json(mcqs: List[Dict], mcq_set: Dict):
    """Export MCQs as JSON"""
    # Convert datetime objects and remove ObjectIds for JSON serialization
    json_mcqs = []
    for mcq in mcqs:
        json_mcq = mcq.copy()
        # Remove MongoDB ObjectId
        if "_id" in json_mcq:
            del json_mcq["_id"]
        # Convert datetime objects to strings
        if "created_at" in json_mcq and json_mcq["created_at"]:
            json_mcq["created_at"] = json_mcq["created_at"].isoformat()
        if "last_reviewed" in json_mcq and json_mcq["last_reviewed"]:
            json_mcq["last_reviewed"] = json_mcq["last_reviewed"].isoformat()
        json_mcqs.append(json_mcq)
    
    # Clean mcq_set data
    clean_mcq_set = mcq_set.copy()
    if "_id" in clean_mcq_set:
        del clean_mcq_set["_id"]
    
    export_data = {
        "set_info": {
            "set_name": clean_mcq_set["set_name"],
            "total_questions": len(json_mcqs),
            "question_type": clean_mcq_set.get("question_type", "single_correct"),
            "difficulty_distribution": clean_mcq_set.get("difficulty_distribution", {}),
            "generated_at": clean_mcq_set["generated_at"].isoformat() if clean_mcq_set.get("generated_at") else None,
            "exported_at": datetime.utcnow().isoformat()
        },
        "mcqs": json_mcqs
    }
    
    # Create response
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    filename = f"{clean_mcq_set['set_name']}_mcqs_{timestamp}.json"
    
    return JSONResponse(
        content=export_data,
        headers={"Content-Disposition": f"attachment; filename={filename}"}
    )

def export_mcqs_pdf(mcqs: List[Dict], mcq_set: Dict):
    """Export MCQs as PDF (basic implementation)"""
    # For now, return a simple text-based PDF export
    # Future enhancement: Use reportlab for better PDF formatting
    
    clean_set_name = mcq_set.get("set_name", "MCQs")
    content = f"MCQ Set: {clean_set_name}\n"
    content += f"Total Questions: {len(mcqs)}\n"
    content += f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n\n"
    content += "=" * 50 + "\n\n"
    
    for i, mcq in enumerate(mcqs, 1):
        content += f"Question {i}: {mcq.get('question_text', '')}\n\n"
        
        options = mcq.get("options", [])
        for j, opt in enumerate(options):
            letter = chr(65 + j)  # A, B, C, D
            marker = " * " if opt.get("is_correct", False) else "   "
            content += f"{marker}{letter}. {opt.get('option_text', '')}\n"
        
        content += f"\nExplanation: {mcq.get('explanation', '')}\n"
        content += f"Topic: {mcq.get('topic', '')}, Difficulty: {mcq.get('difficulty', '')}\n"
        content += "-" * 40 + "\n\n"
    
    # Create response
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    clean_filename = clean_set_name.replace(" ", "_")
    filename = f"{clean_filename}_mcqs_{timestamp}.txt"
    
    from fastapi.responses import Response
    return Response(
        content=content,
        media_type="text/plain",
        headers={"Content-Disposition": f"attachment; filename={filename}"}
    )

# Podcast API Routes
@app.post("/api/generate-transcript")
async def generate_podcast_transcript(request: PodcastTranscriptRequest):
    """Generate podcast transcript from uploaded document using AI (Step 1)"""
    try:
        # Get document info
        doc = await documents_collection.find_one({"file_id": request.file_id})
        if not doc:
            raise HTTPException(status_code=404, detail="Document not found")
        
        file_path = doc["file_path"]
        file_type = doc["file_type"]
        filename = doc["filename"]
        
        # Check if transcript already exists for this document
        existing_transcript = await podcast_transcripts_collection.find_one({"file_id": request.file_id})
        if existing_transcript:
            # Delete existing transcript to regenerate
            await podcast_transcripts_collection.delete_one({"file_id": request.file_id})
        
        # Generate podcast script using AI
        script_data = await generate_podcast_script(
            file_path, 
            file_type, 
            filename,
            request.podcast_length,
            request.content_focus
        )
        
        # Generate transcript ID
        transcript_id = f"transcript_{uuid.uuid4().hex[:8]}"
        
        # Create transcript record
        transcript_data = {
            "transcript_id": transcript_id,
            "title": script_data["title"],
            "description": script_data["description"],
            "script_text": script_data["script_text"],
            "podcast_length": request.podcast_length,
            "estimated_duration": script_data["estimated_duration"],
            "source_doc": filename,
            "file_id": request.file_id,
            "created_at": datetime.utcnow(),
            "is_editable": True
        }
        
        # Save transcript to database
        await podcast_transcripts_collection.insert_one(transcript_data)
        
        # Prepare response data
        response_data = transcript_data.copy()
        if "_id" in response_data:
            del response_data["_id"]
        if "created_at" in response_data:
            response_data["created_at"] = response_data["created_at"].isoformat()
        
        return response_data
            
    except Exception as e:
        logger.error(f"Transcript generation error: {e}")
        raise HTTPException(status_code=500, detail=f"Transcript generation failed: {str(e)}")

@app.post("/api/generate-podcast")
async def generate_podcast_from_transcript(request: PodcastGenerationRequest):
    """Generate podcast audio from transcript using AI voices (Step 2)"""
    try:
        # Get transcript info
        transcript = await podcast_transcripts_collection.find_one({"transcript_id": request.transcript_id})
        if not transcript:
            raise HTTPException(status_code=404, detail="Transcript not found")
        
        # Check if podcast already exists for this transcript
        existing_podcast = await podcasts_collection.find_one({"transcript_id": request.transcript_id})
        if existing_podcast:
            # Delete existing podcast and audio files to regenerate
            if existing_podcast.get("audio_file_path") and os.path.exists(existing_podcast["audio_file_path"]):
                try:
                    os.remove(existing_podcast["audio_file_path"])
                    # Also remove WAV version if exists
                    wav_path = existing_podcast["audio_file_path"].replace(".mp3", ".wav")
                    if os.path.exists(wav_path):
                        os.remove(wav_path)
                except Exception as e:
                    logger.warning(f"Could not delete existing audio file: {e}")
            await podcasts_collection.delete_one({"transcript_id": request.transcript_id})
        
        # Generate podcast ID
        podcast_id = f"podcast_{uuid.uuid4().hex[:8]}"
        
        # Create initial podcast record with processing status
        podcast_data = {
            "podcast_id": podcast_id,
            "transcript_id": request.transcript_id,
            "title": transcript["title"],
            "description": transcript["description"],
            "script_text": transcript["script_text"],
            "audio_file_path": "",
            "duration_seconds": 0,
            "duration_formatted": "00:00",
            "voice_style": request.voice_style,
            "voice_gender": request.voice_gender,
            "voice_accent": request.voice_accent,
            "language": "en",
            "chapters": [],  # Can be enhanced later
            "file_format": "mp3",
            "file_size_bytes": 0,
            "source_doc": transcript["source_doc"],
            "file_id": transcript["file_id"],
            "created_at": datetime.utcnow(),
            "is_processing": True,
            "processing_status": "generating"
        }
        
        # Save initial podcast to database
        await podcasts_collection.insert_one(podcast_data)
        
        try:
            # Generate audio from transcript
            audio_data = await generate_audio_from_script(
                transcript["script_text"],
                request.voice_style,
                request.voice_gender,
                request.voice_accent
            )
            
            # Update podcast with audio information
            await podcasts_collection.update_one(
                {"podcast_id": podcast_id},
                {"$set": {
                    "audio_file_path": audio_data["audio_file_path"],
                    "duration_seconds": audio_data["duration_seconds"],
                    "duration_formatted": audio_data["duration_formatted"],
                    "file_size_bytes": audio_data["file_size_bytes"],
                    "is_processing": False,
                    "processing_status": "completed"
                }}
            )
            
            # Create or update podcast set
            set_id = f"podcast_set_{uuid.uuid4()}"
            existing_set = await podcast_sets_collection.find_one({"file_id": transcript["file_id"]})
            if existing_set:
                set_id = existing_set["set_id"]
                await podcast_sets_collection.update_one(
                    {"file_id": transcript["file_id"]},
                    {"$set": {"last_modified": datetime.utcnow()}}
                )
            else:
                podcast_set = {
                    "set_id": set_id,
                    "set_name": f"{transcript['source_doc'].rsplit('.', 1)[0]} - Podcast",
                    "file_id": transcript["file_id"],
                    "session_id": "podcast_session",
                    "document_title": transcript["source_doc"].rsplit('.', 1)[0],
                    "total_podcasts": 1,
                    "generated_at": datetime.utcnow(),
                    "last_modified": datetime.utcnow(),
                    "is_editable": True
                }
                await podcast_sets_collection.insert_one(podcast_set)
            
            # Prepare response data
            response_data = podcast_data.copy()
            response_data.update(audio_data)
            response_data["set_id"] = set_id
            response_data["is_processing"] = False
            response_data["processing_status"] = "completed"
            
            # Clean response data
            if "_id" in response_data:
                del response_data["_id"]
            if "created_at" in response_data:
                response_data["created_at"] = response_data["created_at"].isoformat()
            
            return response_data
            
        except Exception as audio_error:
            # Update podcast status to failed
            await podcasts_collection.update_one(
                {"podcast_id": podcast_id},
                {"$set": {
                    "is_processing": False,
                    "processing_status": "failed"
                }}
            )
            raise HTTPException(status_code=500, detail=f"Audio generation failed: {str(audio_error)}")
            
    except Exception as e:
        logger.error(f"Podcast generation error: {e}")
        raise HTTPException(status_code=500, detail=f"Podcast generation failed: {str(e)}")

@app.get("/api/transcripts/{file_id}")
async def get_transcript(file_id: str):
    """Get generated transcript for a document"""
    try:
        # Get transcript
        transcript = await podcast_transcripts_collection.find_one({"file_id": file_id})
        if not transcript:
            raise HTTPException(status_code=404, detail="Transcript not found for this document")
        
        # Clean transcript data for JSON response
        clean_transcript = transcript.copy()
        if "_id" in clean_transcript:
            del clean_transcript["_id"]
        if "created_at" in clean_transcript and clean_transcript["created_at"]:
            clean_transcript["created_at"] = clean_transcript["created_at"].isoformat()
        
        return clean_transcript
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to retrieve transcript: {str(e)}")

@app.get("/api/podcasts/{file_id}")
async def get_podcasts(file_id: str):
    """Get generated podcasts for a document"""
    try:
        # Get podcast set
        podcast_set = await podcast_sets_collection.find_one({"file_id": file_id})
        if not podcast_set:
            raise HTTPException(status_code=404, detail="Podcasts not found for this document")
        
        # Get podcasts
        podcasts_cursor = await podcasts_collection.find({"file_id": file_id}).sort("created_at", 1).to_list(length=10)
        
        # Clean podcasts data for JSON response
        clean_podcasts = []
        for podcast in podcasts_cursor:
            clean_podcast = podcast.copy()
            if "_id" in clean_podcast:
                del clean_podcast["_id"]
            if "created_at" in clean_podcast and clean_podcast["created_at"]:
                clean_podcast["created_at"] = clean_podcast["created_at"].isoformat()
            clean_podcasts.append(clean_podcast)
        
        return {
            "set_id": podcast_set["set_id"],
            "set_name": podcast_set["set_name"],
            "podcasts": clean_podcasts,
            "total_podcasts": podcast_set["total_podcasts"],
            "generated_at": podcast_set["generated_at"].isoformat() if podcast_set["generated_at"] else None,
            "last_modified": podcast_set["last_modified"].isoformat() if podcast_set["last_modified"] else None,
            "is_editable": podcast_set.get("is_editable", True)
        }
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to retrieve podcasts: {str(e)}")

@app.put("/api/transcripts/edit")
async def edit_transcript(request: PodcastEditRequest):
    """Edit podcast transcript"""
    try:
        # Find and update the transcript
        result = await podcast_transcripts_collection.update_one(
            {"transcript_id": request.transcript_id},
            {"$set": {
                "title": request.title,
                "script_text": request.script_text
            }}
        )
        
        if result.modified_count == 0:
            raise HTTPException(status_code=404, detail="Transcript not found")
        
        return {"message": "Transcript updated successfully"}
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Transcript update failed: {str(e)}")

@app.post("/api/podcasts/regenerate")
async def regenerate_podcast_audio(request: PodcastRegenerateRequest):
    """Regenerate podcast audio with different voice settings"""
    try:
        # Get transcript
        transcript = await podcast_transcripts_collection.find_one({"transcript_id": request.transcript_id})
        if not transcript:
            raise HTTPException(status_code=404, detail="Transcript not found")
        
        # Find existing podcast
        existing_podcast = await podcasts_collection.find_one({"transcript_id": request.transcript_id})
        if existing_podcast:
            # Clean up old audio file
            if existing_podcast.get("audio_file_path") and os.path.exists(existing_podcast["audio_file_path"]):
                try:
                    os.remove(existing_podcast["audio_file_path"])
                    wav_path = existing_podcast["audio_file_path"].replace(".mp3", ".wav")
                    if os.path.exists(wav_path):
                        os.remove(wav_path)
                except Exception as e:
                    logger.warning(f"Could not delete old audio file: {e}")
        
        # Mark as processing
        await podcasts_collection.update_one(
            {"transcript_id": request.transcript_id},
            {"$set": {
                "is_processing": True,
                "processing_status": "regenerating",
                "voice_style": request.voice_style,
                "voice_gender": request.voice_gender,
                "voice_accent": request.voice_accent
            }}
        )
        
        try:
            # Generate new audio
            audio_data = await generate_audio_from_script(
                transcript["script_text"],
                request.voice_style,
                request.voice_gender,
                request.voice_accent
            )
            
            # Update podcast with new audio
            await podcasts_collection.update_one(
                {"transcript_id": request.transcript_id},
                {"$set": {
                    "audio_file_path": audio_data["audio_file_path"],
                    "duration_seconds": audio_data["duration_seconds"],
                    "duration_formatted": audio_data["duration_formatted"],
                    "file_size_bytes": audio_data["file_size_bytes"],
                    "is_processing": False,
                    "processing_status": "completed"
                }}
            )
            
            return {"message": "Podcast audio regenerated successfully"}
            
        except Exception as audio_error:
            # Mark as failed
            await podcasts_collection.update_one(
                {"transcript_id": request.transcript_id},
                {"$set": {
                    "is_processing": False,
                    "processing_status": "failed"
                }}
            )
            raise HTTPException(status_code=500, detail=f"Audio regeneration failed: {str(audio_error)}")
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Regeneration failed: {str(e)}")

@app.delete("/api/podcasts/{podcast_id}")
async def delete_podcast(podcast_id: str):
    """Delete a specific podcast"""
    try:
        # Get podcast info before deletion
        podcast = await podcasts_collection.find_one({"podcast_id": podcast_id})
        if not podcast:
            raise HTTPException(status_code=404, detail="Podcast not found")
        
        # Clean up audio files
        if podcast.get("audio_file_path") and os.path.exists(podcast["audio_file_path"]):
            try:
                os.remove(podcast["audio_file_path"])
                # Also remove WAV version
                wav_path = podcast["audio_file_path"].replace(".mp3", ".wav")
                if os.path.exists(wav_path):
                    os.remove(wav_path)
            except Exception as e:
                logger.warning(f"Could not delete audio file: {e}")
        
        # Delete the podcast
        await podcasts_collection.delete_one({"podcast_id": podcast_id})
        
        # Update set total count
        remaining_count = await podcasts_collection.count_documents({"file_id": podcast["file_id"]})
        await podcast_sets_collection.update_one(
            {"file_id": podcast["file_id"]},
            {"$set": {
                "total_podcasts": remaining_count,
                "last_modified": datetime.utcnow()
            }}
        )
        
        return {"message": "Podcast deleted successfully"}
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Podcast deletion failed: {str(e)}")

@app.get("/api/podcasts/audio/{podcast_id}")
async def get_podcast_audio(podcast_id: str, format: str = "mp3"):
    """Stream or download podcast audio file"""
    try:
        logger.info(f"Getting podcast audio for ID: {podcast_id}, format: {format}")
        
        # Get podcast info
        podcast = await podcasts_collection.find_one({"podcast_id": podcast_id})
        if not podcast:
            logger.error(f"Podcast not found: {podcast_id}")
            raise HTTPException(status_code=404, detail="Podcast not found")
        
        audio_file_path = podcast.get("audio_file_path")
        logger.info(f"Audio file path: {audio_file_path}")
        
        if not audio_file_path:
            logger.error(f"No audio file path found for podcast: {podcast_id}")
            raise HTTPException(status_code=404, detail="Audio file path not found")
            
        if not os.path.exists(audio_file_path):
            logger.error(f"Audio file does not exist: {audio_file_path}")
            raise HTTPException(status_code=404, detail="Audio file not found on disk")
        
        # Handle format conversion if needed
        if format.lower() == "wav" and audio_file_path.endswith(".mp3"):
            wav_path = audio_file_path.replace(".mp3", ".wav")
            if not os.path.exists(wav_path):
                logger.info(f"Converting {audio_file_path} to WAV format")
                # Convert to WAV
                wav_path = await convert_audio_format(audio_file_path, "wav")
            audio_file_path = wav_path
        
        # Verify final file exists
        if not os.path.exists(audio_file_path):
            logger.error(f"Final audio file does not exist: {audio_file_path}")
            raise HTTPException(status_code=404, detail="Audio file not accessible")
        
        logger.info(f"Serving audio file: {audio_file_path}")
        
        # Return file
        from fastapi.responses import FileResponse
        media_type = "audio/mpeg" if format.lower() == "mp3" else "audio/wav"
        filename = f"{podcast.get('title', 'podcast').replace(' ', '_')}.{format.lower()}"
        
        return FileResponse(
            audio_file_path,
            media_type=media_type,
            filename=filename
        )
    
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Unexpected error in get_podcast_audio: {str(e)}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"Failed to serve audio file: {str(e)}")

@app.post("/api/podcasts/export")
async def export_podcast(request: PodcastExportRequest):
    """Export podcast in specified format"""
    try:
        logger.info(f"Exporting podcast: {request.podcast_id} in format: {request.export_format}")
        
        # Get podcast info
        podcast = await podcasts_collection.find_one({"podcast_id": request.podcast_id})
        if not podcast:
            logger.error(f"Podcast not found: {request.podcast_id}")
            raise HTTPException(status_code=404, detail="Podcast not found")
        
        audio_file_path = podcast.get("audio_file_path")
        logger.info(f"Audio file path: {audio_file_path}")
        
        if not audio_file_path:
            logger.error(f"No audio file path found for podcast: {request.podcast_id}")
            raise HTTPException(status_code=404, detail="Audio file path not found")
            
        if not os.path.exists(audio_file_path):
            logger.error(f"Audio file does not exist: {audio_file_path}")
            raise HTTPException(status_code=404, detail="Audio file not found on disk")
        
        # Handle format conversion if needed
        export_format = request.export_format.lower()
        if export_format == "wav" and audio_file_path.endswith(".mp3"):
            # Convert to WAV
            wav_path = audio_file_path.replace(".mp3", ".wav")
            if not os.path.exists(wav_path):
                logger.info(f"Converting {audio_file_path} to WAV format")
                wav_path = await convert_audio_format(audio_file_path, "wav")
            export_path = wav_path
        else:
            export_path = audio_file_path
        
        # Verify final export file exists
        if not os.path.exists(export_path):
            logger.error(f"Export file does not exist: {export_path}")
            raise HTTPException(status_code=404, detail="Export file not accessible")
        
        logger.info(f"Exporting file: {export_path}")
        
        # Return file for download
        from fastapi.responses import FileResponse
        media_type = "audio/mpeg" if export_format == "mp3" else "audio/wav"
        safe_title = podcast.get("title", "podcast").replace(" ", "_").replace("/", "_")
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        filename = f"{safe_title}_{timestamp}.{export_format}"
        
        logger.info(f"Serving file with filename: {filename}")
        
        return FileResponse(
            export_path,
            media_type=media_type,
            headers={
                "Content-Disposition": f"attachment; filename={filename}",
                "Cache-Control": "no-cache"
            }
        )
    
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Unexpected error in export_podcast: {str(e)}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"Export failed: {str(e)}")

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8001)