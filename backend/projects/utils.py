# projects/utils.py
import os
import json
import boto3
import openai
import PyPDF2
from docx import Document
from pptx import Presentation
from django.conf import settings

# Configure OpenAI client
openai.api_key = settings.OPENAI_API_KEY

def download_file_from_s3(s3_key):
    """Downloads a file from S3 to a temporary local path."""

    if s3_key.startswith('http'):
        s3_key = s3_key.split('.com/',1)[1]
    s3_client = boto3.client(
        's3',
        aws_access_key_id=settings.AWS_ACCESS_KEY_ID,
        aws_secret_access_key=settings.AWS_SECRET_ACCESS_KEY
    )
    # Ensure the /tmp directory exists
    temp_dir = "/tmp/downloads/"
    os.makedirs(temp_dir, exist_ok=True)
    
    local_path = os.path.join(temp_dir, os.path.basename(s3_key))
    s3_client.download_file(settings.AWS_STORAGE_BUCKET_NAME, s3_key, local_path)
    return local_path

def extract_text_from_file(file_path):
    """Extracts text content from a PDF or DOCX file."""
    text = ""
    try:
        if file_path.endswith('.pdf'):
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    text += page.extract_text() or ""
        elif file_path.endswith('.docx'):
            doc = Document(file_path)
            for para in doc.paragraphs:
                text += para.text + "\n"
        elif file_path.endswith('.txt'):
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
    finally:
        # Clean up the downloaded file
        if os.path.exists(file_path):
            os.remove(file_path)
    return text


def generate_image_for_slide(slide_title, slide_content):
    """Generates an image using DALL-E and returns a file-like object."""
    try:
        prompt = f"""
        A professional, clean, and modern illustration for a presentation slide.
        The slide is titled '{slide_title}' and discusses '{' '.join(slide_content[:2])}'.
        The style should be minimalist and suitable for a business or educational setting. No text in the image.
        """
        
        response = openai.images.generate(
            model="dall-e-3",
            prompt=prompt,
            size="1024x1024",
            quality="standard",
            n=1,
            response_format="b64_json" # Get the image as a base64 string
        )
        
        image_data = base64.b64decode(response.data[0].b64_json)
        #Cost of image generation
        cost = settings.OPENAI_PRICING.get("dall-e-3").get("standard_1024x1024")
        return BytesIO(image_data),cost # Return as an in-memory file
        
    except Exception as e:
        print(f"DALL-E image generation failed: {e}")
        return None,0

def generate_ppt_from_text(text_content,options):
    """
    Uses OpenAI to get slide content and python-pptx to create a presentation.
    """
    slide_count = options.get("slide_count", 8)
    include_images = options.get("include_images", False)
    prompt = f"""
    Based on the following text, create content for EXACTLY {slide_count} presentation slides.
    Return a valid JSON object only, with a single key "slides" which is a list of objects.
    Each object must have three keys: "title" (string), "content" (a list of 4-5 strings for bullet points), and "speaker_notes" (a detailed paragraph).
    CRITICAL: All content must be specific to the uploaded document.
    TEXT: ---
    {text_content[:8000]}
    ---
    """
    
    # Make the actual API call to OpenAI
    response = openai.chat.completions.create(
        model="gpt-5-nano",
        messages=[{"role": "user", "content": prompt}],
        response_format={"type": "json_object"}
    )
    
    # Extract and parse the JSON response
    slide_data_json = json.loads(response.choices[0].message.content)
    slides_data = slide_data_json.get("slides", [])
    text_usage = response.usage
    total_cost = calculate_cost("gpt-5-nano",text_usage) #start with text cost

    prs = Presentation()
    image_slide_layout = prs.slide_layouts[8]
    text_only_slide_layout = prs.slide_layouts[1]
    for slide_info in slides_data:
        image_stream = None
        if include_images:
            print("Generating image for slide:", slide_info.get("title"))
            image_stream,image_cost = generate_image_for_slide(slide_info.get("title"), slide_info.get("content", []))
            total_cost += image_cost


        
        if image_stream:
            slide = prs.slides.add_slide(image_slide_layout)
            # Add text to the main placeholder
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.text = slide_info.get("title", "Untitled Slide")
            for point in slide_info.get("content", []):
                p = tf.add_paragraph()
                p.text = point
                p.level = 1
            
            # Add the image to the image placeholder
            image_placeholder = slide.placeholders[2]
            image_placeholder.insert_picture(image_stream)
        else:
            # Fallback to text-only layout if no image is generated
            slide = prs.slides.add_slide(text_only_slide_layout)
            title = slide.shapes.title
            body = slide.placeholders[1]
            title.text = slide_info.get("title", "Untitled Slide")
            tf = body.text_frame
            tf.clear()
            for point in slide_info.get("content", []):
                p = tf.add_paragraph()
                p.text = point
                p.level = 1
        
        # Add speaker notes
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = slide_info.get("speaker_notes", "")

    # Save the presentation to a temporary file
    temp_dir = "/tmp/presentations/"
    os.makedirs(temp_dir, exist_ok=True)
    output_path = os.path.join(temp_dir, "presentation.pptx")
    prs.save(output_path)
    
    return output_path ,total_cost





def generate_flashcards_from_text(text_content,options):
    """
    Uses OpenAI to generate flashcards as structured JSON data.
    """
    card_count = options.get("cards_count", 20)
    card_type = options.get("card_type", "qa")
    difficulty = options.get("difficulty", "mixed")
    prompt = f"""
    Based on the following text, generate EXACTLY  {card_count} flashcards. The card type should be {card_type} and the difficulty should be {difficulty}.
    Return a valid JSON object only, with a single key "flashcards" which is a list of objects.
    Each object must have three keys: "question" (string), "answer" (a concise string), "topic" (a relevant keyword), "difficulty" (difficulty level of question, can be easy, medium, hard).
    TEXT: ---
    {text_content[:8000]}
    ---
    """

    # Make the actual API call to OpenAI
    response = openai.chat.completions.create(
        model="gpt-5-nano",
        messages=[{"role": "user", "content": prompt}],
        response_format={"type": "json_object"}
    )


    # Extract and parse the JSON response
    flashcard_data_json = json.loads(response.choices[0].message.content)
    usage = response.usage
    total_cost = calculate_cost("gpt-5-nano",usage) #start with text cost
    temp_dir = "/tmp/flashcards/"
    os.makedirs(temp_dir, exist_ok=True)
    output_path = os.path.join(temp_dir, "flashcards.json")
    with open(output_path, 'w') as f:
        json.dump(flashcard_data_json, f, indent=4)
        
    return output_path ,total_cost



def generate_mcqs_from_text(text_content,options):
    """
    Uses OpenAI to generate MCQs as structured JSON data.
    """
    questions_count = options.get("questions_count", 15)
    questions_type = options.get("questions_type", "single_correct")
    difficulty = options.get("difficulty", "mixed")
    prompt = f"""
Based on the following text, generate EXACTLY {questions_count} MCQs. The question type should be {questions_type} and the difficulty should be {difficulty}.

Return a valid JSON object only, with a single key "mcqs", which is a list of question objects.

Each object in the "mcqs" list must have the following fields:
1. "question_text" (string): The question itself.
2. "options" (list of 4 objects): Each object must have:
   - "option_text" (string): The option text.
   - "is_correct" (boolean): Whether this is the correct answer.
   Exactly one option must have "is_correct": true.
3. "explanation" (string): A brief explanation of why the correct option is correct.
4. "difficulty" (string): One of "easy", "medium", or "hard".
5. "bloom_level" (string): The Bloom's taxonomy level of the question. One of:
   - "Remember", "Understand", "Apply", "Analyze", "Evaluate", "Create"
6. "topic" (string): A brief topic or concept this question relates to.

TEXT: ---
{text_content[:8000]}
---
"""


    # Make the actual API call to OpenAI
    response = openai.chat.completions.create(
        model="gpt-5-nano",
        messages=[{"role": "user", "content": prompt}],
        response_format={"type": "json_object"}
    )

    # Extract and parse the JSON response
    mcq_data_json = json.loads(response.choices[0].message.content)
    usage = response.usage
    total_cost = calculate_cost("gpt-5-nano",usage) #start with text cost

    temp_dir = "/tmp/mcqs/"
    os.makedirs(temp_dir, exist_ok=True)
    output_path = os.path.join(temp_dir, "mcqs.json")
    with open(output_path, 'w') as f:
        json.dump(mcq_data_json, f, indent=4)
        
    return output_path ,total_cost



import edge_tts
import asyncio


#old podcast functions
async def generate_audio_from_script(script_text, output_path, voice_name):
    """Uses edge-tts to convert text to an MP3 file."""
    communicate = edge_tts.Communicate(script_text, voice_name)
    await communicate.save(output_path)

# def generate_podcast_from_text(text_content,options):
#     """
#     Uses OpenAI to generate a podcast script and then converts it to an audio file.
#     """
#     # Step 1: Generate the script using OpenAI
#     podcast_length = options.get("podcast_length", "medium")
#     content_focus = options.get("content_focus", "full_document")
#     voice_accent = options.get("voice_accent", "american")
#     voice_gender = options.get("voice_gender", "female")
#     voice_style = options.get("voice_style", "neutral")

#     voice_map = {
#             # American English
#             ("american", "female", "neutral"): "en-US-AriaNeural",
#             ("american", "female", "enthusiastic"): "en-US-JennyNeural", 
#             ("american", "female", "formal"): "en-US-SaraNeural",
#             ("american", "female", "conversational"): "en-US-AriaNeural",
#             ("american", "male", "neutral"): "en-US-GuyNeural",
#             ("american", "male", "enthusiastic"): "en-US-BrianNeural",
#             ("american", "male", "formal"): "en-US-DavisNeural", 
#             ("american", "male", "conversational"): "en-US-GuyNeural",
            
#             # British English
#             ("british", "female", "neutral"): "en-GB-SoniaNeural",
#             ("british", "female", "enthusiastic"): "en-GB-LibbyNeural", 
#             ("british", "female", "formal"): "en-GB-SoniaNeural",
#             ("british", "female", "conversational"): "en-GB-MaisieNeural",
#             ("british", "male", "neutral"): "en-GB-RyanNeural",
#             ("british", "male", "enthusiastic"): "en-GB-ThomasNeural",
#             ("british", "male", "formal"): "en-GB-RyanNeural", 
#             ("british", "male", "conversational"): "en-GB-AlfieNeural",
            
#             # Indian English
#             ("indian", "female", "neutral"): "en-IN-NeerjaNeural",
#             ("indian", "female", "enthusiastic"): "en-IN-NeerjaNeural", 
#             ("indian", "female", "formal"): "en-IN-NeerjaNeural",
#             ("indian", "female", "conversational"): "en-IN-NeerjaNeural",
#             ("indian", "male", "neutral"): "en-IN-PrabhatNeural",
#             ("indian", "male", "enthusiastic"): "en-IN-PrabhatNeural",
#             ("indian", "male", "formal"): "en-IN-PrabhatNeural", 
#             ("indian", "male", "conversational"): "en-IN-PrabhatNeural",
            
#             # Australian English
#             ("australian", "female", "neutral"): "en-AU-NatashaNeural",
#             ("australian", "female", "enthusiastic"): "en-AU-NatashaNeural", 
#             ("australian", "female", "formal"): "en-AU-NatashaNeural",
#             ("australian", "female", "conversational"): "en-AU-NatashaNeural",
#             ("australian", "male", "neutral"): "en-AU-WilliamNeural",
#             ("australian", "male", "enthusiastic"): "en-AU-WilliamNeural",
#             ("australian", "male", "formal"): "en-AU-WilliamNeural", 
#             ("australian", "male", "conversational"): "en-AU-WilliamNeural",
            
#             # Canadian English
#             ("canadian", "female", "neutral"): "en-CA-ClaraNeural",
#             ("canadian", "female", "enthusiastic"): "en-CA-ClaraNeural", 
#             ("canadian", "female", "formal"): "en-CA-ClaraNeural",
#             ("canadian", "female", "conversational"): "en-CA-ClaraNeural",
#             ("canadian", "male", "neutral"): "en-CA-LiamNeural",
#             ("canadian", "male", "enthusiastic"): "en-CA-LiamNeural",
#             ("canadian", "male", "formal"): "en-CA-LiamNeural", 
#             ("canadian", "male", "conversational"): "en-CA-LiamNeural"
#         }
    

#     voice_name = voice_map.get((voice_accent, voice_gender, voice_style), "en-US-AriaNeural")
#     script_prompt = f"""
#     Based on the following text, write podcast script with a target lenght of {podcast_length} that focuses on {content_focus}.
#     The tone should be engaging and conversational.
#     Return a valid JSON object only, with a single key "script" which is an object.
#     The script object must have three keys: "title" (string), "description" (string), and "body" (a long string with the full podcast script).
#     TEXT: ---
#     {text_content[:6000]}
#     ---
#     """
#     response = openai.chat.completions.create(
#         model="gpt-5-nano",
#         messages=[{"role": "user", "content": script_prompt}],
#         response_format={"type": "json_object"}
#     )
#     script_data = json.loads(response.choices[0].message.content).get("script", {})
#     usage = response.usage
#     script_text = script_data.get("body", "Script could not be generated.")

#     # Step 2: Convert the script text to an audio file
#     temp_dir = "/tmp/podcasts/"
#     os.makedirs(temp_dir, exist_ok=True)
#     output_path = os.path.join(temp_dir, "podcast.mp3")

#     # Run the async TTS function
#     asyncio.run(generate_audio_from_script(script_text, output_path,voice_name))

#     # Return the path to the generated MP3 and the script data
#     return output_path, script_data , usage

#new podcast functions

def generate_podcast_script_from_text(text_content,options,document_title):
    """
    Uses OpenAI to generate a podcast script.
    """
    podcast_length = options.get('podcast_length', 'medium')
    content_focus = options.get('content_focus', 'full_document')
    
    length_instruction = ""
    if podcast_length == "quick":
            length_instruction = "Create a 2-3 minute podcast script focusing on key takeaways and essential points."
    elif podcast_length == "comprehensive":
            length_instruction = "Create a 7-10 minute podcast script with detailed explanations and comprehensive coverage."
    else:  # medium
            length_instruction = "Create a 4-6 minute podcast script balancing key concepts with engaging explanations."
            
    focus_instruction = ""
    if content_focus == "key_concepts":
            focus_instruction = "Focus primarily on the most important concepts, definitions, and core ideas."
    elif content_focus == "summary":
            focus_instruction = "Provide a comprehensive summary hitting all major points concisely."
    else:  # full_document
            focus_instruction = "Cover the full document content in an engaging, structured manner."

    
    script_prompt = f"""
        Please analyze this document and create an engaging podcast script based on the content.
        
        {length_instruction}
        {focus_instruction}
        
        ***Return a valid JSON object only, with a single key "script" which is an object.***
        ***The script object must have three keys: "title" (string), "description" (string), and "body" (a long string with the full podcast script).***
        
        SCRIPT WRITING GUIDELINES:
        1. Write in a conversational, engaging tone as if speaking directly to the listener
        2. Use natural speech patterns with appropriate pauses (indicated by commas and periods)
        3. Include smooth transitions between topics
        4. Explain complex concepts in simple, accessible language
        5. Use rhetorical questions to engage the listener
        6. Include brief examples or analogies when helpful
        7. Maintain an educational yet entertaining style
        8. Structure the content logically with clear flow
        9. Include verbal signposts like "First," "Next," "Finally," etc.
        10. End with actionable takeaways or thought-provoking questions
        
        INTRO TEMPLATE:
        "Welcome to this learning session on [topic]. I'm excited to share some fascinating insights about [subject] that will help you understand [key benefit]. Let's dive right in!"
        
        CONCLUSION TEMPLATE:
        "To wrap up today's session, we've covered [key points]. The main takeaway is [core message]. I hope this helps you [practical application]. Thanks for listening, and keep learning!"
        
        Document title: {document_title}
        Document content: {text_content[:6000]}...
        
        Create a script that transforms this written content into an engaging audio learning experience.
        """
        
    response = openai.chat.completions.create(
        model="gpt-5-nano",
        messages=[{"role": "user", "content": script_prompt}],
        response_format={"type": "json_object"}
    )
    script_data = json.loads(response.choices[0].message.content).get("script", {})
    usage = response.usage
    # total_cost = calculate_cost("gpt-5-nano",usage) #start with text cost
    # script_text = script_data.get("body", "Script could not be generated.")
    return script_data, usage

def generate_podcast_audio_from_script(script_data,options):
    """
    Uses OpenAI to generate a podcast from a script.
    """
    voice_style = options.get('voice_style', 'neutral')
    voice_gender = options.get('voice_gender', 'female')
    voice_accent = options.get('voice_accent', 'american')

    voice_map = {
            # American English
            ("american", "female", "neutral"): "en-US-AriaNeural",
            ("american", "female", "enthusiastic"): "en-US-JennyNeural", 
            ("american", "female", "formal"): "en-US-SaraNeural",
            ("american", "female", "conversational"): "en-US-AriaNeural",
            ("american", "male", "neutral"): "en-US-GuyNeural",
            ("american", "male", "enthusiastic"): "en-US-BrianNeural",
            ("american", "male", "formal"): "en-US-DavisNeural", 
            ("american", "male", "conversational"): "en-US-GuyNeural",
            
            # British English
            ("british", "female", "neutral"): "en-GB-SoniaNeural",
            ("british", "female", "enthusiastic"): "en-GB-LibbyNeural", 
            ("british", "female", "formal"): "en-GB-SoniaNeural",
            ("british", "female", "conversational"): "en-GB-MaisieNeural",
            ("british", "male", "neutral"): "en-GB-RyanNeural",
            ("british", "male", "enthusiastic"): "en-GB-ThomasNeural",
            ("british", "male", "formal"): "en-GB-RyanNeural", 
            ("british", "male", "conversational"): "en-GB-AlfieNeural",
            
            # Indian English
            ("indian", "female", "neutral"): "en-IN-NeerjaNeural",
            ("indian", "female", "enthusiastic"): "en-IN-NeerjaNeural", 
            ("indian", "female", "formal"): "en-IN-NeerjaNeural",
            ("indian", "female", "conversational"): "en-IN-NeerjaNeural",
            ("indian", "male", "neutral"): "en-IN-PrabhatNeural",
            ("indian", "male", "enthusiastic"): "en-IN-PrabhatNeural",
            ("indian", "male", "formal"): "en-IN-PrabhatNeural", 
            ("indian", "male", "conversational"): "en-IN-PrabhatNeural",
            
            # Australian English
            ("australian", "female", "neutral"): "en-AU-NatashaNeural",
            ("australian", "female", "enthusiastic"): "en-AU-NatashaNeural", 
            ("australian", "female", "formal"): "en-AU-NatashaNeural",
            ("australian", "female", "conversational"): "en-AU-NatashaNeural",
            ("australian", "male", "neutral"): "en-AU-WilliamNeural",
            ("australian", "male", "enthusiastic"): "en-AU-WilliamNeural",
            ("australian", "male", "formal"): "en-AU-WilliamNeural", 
            ("australian", "male", "conversational"): "en-AU-WilliamNeural",
            
            # Canadian English
            ("canadian", "female", "neutral"): "en-CA-ClaraNeural",
            ("canadian", "female", "enthusiastic"): "en-CA-ClaraNeural", 
            ("canadian", "female", "formal"): "en-CA-ClaraNeural",
            ("canadian", "female", "conversational"): "en-CA-ClaraNeural",
            ("canadian", "male", "neutral"): "en-CA-LiamNeural",
            ("canadian", "male", "enthusiastic"): "en-CA-LiamNeural",
            ("canadian", "male", "formal"): "en-CA-LiamNeural", 
            ("canadian", "male", "conversational"): "en-CA-LiamNeural"
        }
    

    voice_name = voice_map.get((voice_accent, voice_gender, voice_style), "en-US-AriaNeural")
    temp_dir = "/tmp/podcasts/"
    os.makedirs(temp_dir, exist_ok=True)
    output_path = os.path.join(temp_dir, "podcast.mp3")
    asyncio.run(generate_audio_from_script(script_data, output_path,voice_name))
    return output_path
    

    


# projects/utils.py

def calculate_cost(model_name, usage):
    print("usage",usage)
    # logger.info("usage",usage)
    """Calculates the cost of an OpenAI API call."""
    pricing = settings.OPENAI_PRICING.get(model_name)
    if not pricing:
        return 0.0

    prompt_tokens = usage.prompt_tokens
    completion_tokens = usage.completion_tokens

    cost = (prompt_tokens * pricing['input']) + (completion_tokens * pricing['output'])
    return cost



